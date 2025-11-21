import os
import tempfile
import zipfile
from pathlib import Path
import streamlit as st
from pptx import Presentation
import fitz  # PyMuPDF
from PIL import Image
from openai import AzureOpenAI
from dotenv import load_dotenv
import re
from pptx.dml.color import RGBColor



# =========================
# ページ設定
# =========================
st.set_page_config(page_title="企画プラットフォーム", layout="wide")
st.title("企画プラットフォーム")

# =========================
# セッション初期化
# =========================
if "selected_mode" not in st.session_state:
    st.session_state["selected_mode"] = None
if "message_center" not in st.session_state:
    st.session_state["message_center"] = ""
if "message_right" not in st.session_state:
    st.session_state["message_right"] = ""
if "uploaded_docs" not in st.session_state:
    st.session_state["uploaded_docs"] = []
if "pptx_path" not in st.session_state:
    st.session_state["pptx_path"] = None
if "edited_texts" not in st.session_state:
    st.session_state["edited_texts"] = {}
if "orien_outline_text" not in st.session_state:
    st.session_state["orien_outline_text"] = ""
if "orien_company_text" not in st.session_state:
    st.session_state["orien_company_text"] = ""
if "final_pptx_path" not in st.session_state:
    st.session_state["final_pptx_path"] = None



# =========================
# Azure OpenAI 設定
# =========================
load_dotenv()
client = AzureOpenAI(
    api_key=os.getenv("OPENAI_API_KEY"),
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
    api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
)
DEPLOYMENT = os.getenv("AZURE_OPENAI_DEPLOYMENT", "gpt-4o")


# =========================
# 古いセッションの自動クリーンアップ
# =========================
def cleanup_old_sessions(days: int = 1):
    """
    最終アクセスから days 日以上経過したセッションディレクトリを削除
    """
    if not BASE_ROOT.exists():
        return

    now = datetime.now()
    for child in BASE_ROOT.iterdir():
        if not child.is_dir():
            continue

        marker = child / ".last_access"
        try:
            if marker.exists():
                ts = datetime.fromisoformat(marker.read_text(encoding="utf-8"))
            else:
                # マーカーがない場合はディレクトリの更新時刻を使う
                ts = datetime.fromtimestamp(child.stat().st_mtime)

            if now - ts > timedelta(days=days):
                shutil.rmtree(child, ignore_errors=True)
        except Exception:
            # 読み取り・削除で何かあってもアプリを落とさない
            continue



# =========================
# セッション専用ディレクトリのヘルパーを作る
# =========================
import uuid
import shutil
from datetime import datetime, timedelta

BASE_ROOT = Path("/home/streamlit_workspace")


def get_session_dir() -> Path:

    cleanup_old_sessions(days=1) 
    """
    セッションごとに一意の作業ディレクトリを返す。
    例）/home/streamlit_workspace/20250201_120000_ab12cd34/
    """
    if "session_id" not in st.session_state:
        sid = datetime.now().strftime("%Y%m%d_%H%M%S") + "_" + uuid.uuid4().hex[:8]
        st.session_state["session_id"] = sid

    session_dir = BASE_ROOT / st.session_state["session_id"]
    session_dir.mkdir(parents=True, exist_ok=True)

    # 最終アクセス時刻を記録しておく（自動クリーンアップ用）
    (session_dir / ".last_access").write_text(datetime.now().isoformat(), encoding="utf-8")

    return session_dir


# =========================
# ファイル読込関数
# =========================
def read_txt(path):
    for enc in ("utf-8", "utf-8-sig", "cp932"):
        try:
            with open(path, "r", encoding=enc, errors="ignore") as f:
                return f.read()
        except Exception:
            continue
    return ""

def read_pdf(path):
    try:
        doc = fitz.open(path)
        return "\n".join(page.get_text("text") for page in doc)
    except Exception:
        return ""

def read_pptx_text(path):
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shp in slide.shapes:
                if hasattr(shp, "text") and shp.text:
                    texts.append(shp.text)
        return "\n".join(texts)
    except Exception:
        return ""

# =========================
# PPT → 画像変換関数
# =========================
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont
import io, os

def pptx_to_images(pptx_path: Path) -> list[Image.Image]:
    """
    PowerPointファイルをスライドレイアウト通りに簡易描画して画像リストで返す。
    - 日本語フォント対応
    - テキスト・画像を元の位置(left, top, width, height)に再配置
    """
    images: list[Image.Image] = []

    # ---- 日本語フォント設定 ----
    FONT_CANDIDATES = [
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/fonts-japanese-gothic.ttf",
        "C:/Windows/Fonts/meiryo.ttc",
        "/System/Library/Fonts/ヒラギノ角ゴシック W4.ttc",
        "/System/Library/Fonts/Helvetica.ttc",
    ]
    font_path = next((f for f in FONT_CANDIDATES if os.path.exists(f)), None)
    if font_path:
        font_small = ImageFont.truetype(font_path, 20)
    else:
        font_small = ImageFont.load_default()

    try:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            # スライドサイズ（EMU → px換算）
            width_px = int(prs.slide_width / 9525)
            height_px = int(prs.slide_height / 9525)

            # 白背景キャンバス
            img = Image.new("RGB", (width_px, height_px), "white")
            draw = ImageDraw.Draw(img)

            # === スライド上の図形を順に描画 ===
            for shp in slide.shapes:
                left = int(shp.left / 9525)
                top = int(shp.top / 9525)
                width = int(shp.width / 9525)
                height = int(shp.height / 9525)

                # 図形タイプで分岐
                stype = shp.shape_type

                # 画像
                if stype == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_bytes = io.BytesIO(shp.image.blob)
                        pic = Image.open(image_bytes).convert("RGB")
                        pic = pic.resize((width, height))
                        img.paste(pic, (left, top))
                    except Exception:
                        draw.rectangle([left, top, left + width, top + height], outline="gray")
                        draw.text((left + 4, top + 4), "画像読み込み失敗", font=font_small, fill="red")

                # テキスト付き図形
                elif getattr(shp, "has_text_frame", False):
                    text = shp.text.strip()
                    if text:
                        # テキスト枠（背景塗り）
                        draw.rectangle([left, top, left + width, top + height], outline="lightgray", fill=None)
                        # テキスト（簡易左寄せ）
                        lines = text.replace("\r", "").split("\n")
                        y = top + 5
                        for line in lines:
                            draw.text((left + 8, y), line[:40], font=font_small, fill="black")
                            y += 24

                # 図形（塗りつぶしのみ）
                else:
                    draw.rectangle([left, top, left + width, top + height], outline="lightgray", fill=None)

            # スライド番号
            draw.text((20, height_px - 40), f"Slide {i+1}", font=font_small, fill="gray")

            images.append(img)

        return images

    except Exception as e:
        st.error(f"PPT変換エラー: {e}")
        return []
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import base64

def emu_to_percent(val_emu, total_emu):
    """EMU単位をスライド全体に対する%へ変換"""
    try:
        return float(val_emu) / float(total_emu) * 100.0
    except Exception:
        return 0.0


def color_to_css(rgb):
    """RGBColor → CSSカラーコード"""
    if not rgb:
        return None
    if isinstance(rgb, RGBColor):
        return f"#{rgb.rgb:06X}"
    try:
        return f"#{int(rgb):06X}"
    except Exception:
        return None


def extract_slide_model(prs, slide_index=0):
    """
    PowerPointスライド内の図形を走査し、
    Streamlit用のHTML描画モデルに変換する。
    - PICTURE：画像
    - TEXT：テキストボックス
    - TABLE：セルの文字を連結して1つのテキストブロックとして描画（★追加）
    """
    slide = prs.slides[slide_index]
    sw, sh = prs.slide_width, prs.slide_height
    blocks = []

    def add_block(shape, offset_left=0, offset_top=0):
        stype = shape.shape_type
        name = getattr(shape, "name", "")
        editable = name.startswith("Edit_") or name.startswith("EDIT_")

        left = shape.left + offset_left
        top = shape.top + offset_top
        width = shape.width
        height = shape.height

        base = {
            "name": name,
            "editable": editable,
            "left": emu_to_percent(left, sw),
            "top": emu_to_percent(top, sh),
            "width": emu_to_percent(width, sw),
            "height": emu_to_percent(height, sh),
        }

        # グループ処理
        if stype == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                add_block(child, offset_left=left, offset_top=top)
            return

        # 画像処理
        if stype == MSO_SHAPE_TYPE.PICTURE:
            try:
                content_type = getattr(shape.image, "content_type", "image/png")
                b64 = base64.b64encode(shape.image.blob).decode("ascii")
                base["type"] = "picture"
                base["src"] = f"data:{content_type};base64,{b64}"
            except Exception:
                pass
            blocks.append(base)
            return

        # 塗り（背景色）
        fill_css = None
        try:
            if shape.fill and shape.fill.type == 1:  # solid fill
                fill_css = color_to_css(shape.fill.fore_color.rgb)
        except Exception:
            pass

        # ★ TABLE（表）の処理を追加：セルのテキストを連結して1ブロックとして描画
        if stype == MSO_SHAPE_TYPE.TABLE:
            try:
                table = shape.table
                rows_text = []
                for row in table.rows:
                    cells = [
                        cell.text.replace("\n", " ").strip()
                        for cell in row.cells
                    ]
                    # 全部空なら無視
                    if any(cells):
                        rows_text.append(" | ".join(cells))
                text = "\n".join(rows_text).strip()
                if text:
                    blocks.append(
                        {
                            **base,
                            "type": "text",
                            "text": text,
                            "fill": fill_css,
                        }
                    )
            except Exception:
                # 失敗したら単なるボックスとして描画
                blocks.append({**base, "type": "box", "fill": fill_css})
            return

        # テキスト付き図形
        if getattr(shape, "has_text_frame", False):
            blocks.append({**base, "type": "text", "text": shape.text, "fill": fill_css})
        else:
            # 図形（塗りだけ）
            blocks.append({**base, "type": "box", "fill": fill_css})

    for shape in slide.shapes:
        add_block(shape, 0, 0)

    return {"blocks": blocks}


def render_slide_html(model, edited_texts):
    """
    extract_slide_model()で抽出した構造をもとに、
    Streamlit内でスライドの見た目を再現するHTMLを生成。
    """
    blocks = model["blocks"]

    html = """
    <div style="position:relative; width:100%; padding-top:56.25%; background:#f8f9fb;
                border-radius:14px; box-shadow:0 4px 16px rgba(0,0,0,0.08); overflow:hidden;">
      <div style="position:absolute; inset:0; background:white;">
    """

    for b in blocks:
        style = (
            f"position:absolute; left:{b['left']}%; top:{b['top']}%; "
            f"width:{b['width']}%; height:{b['height']}%;"
        )
        content = ""

        # 画像
        if b.get("type") == "picture" and b.get("src"):
            content = (
                f'<img src="{b["src"]}" style="width:100%;height:100%;object-fit:contain;">'
            )

        # テキスト
        elif b.get("type") == "text":
            text_val = edited_texts.get(b["name"], b.get("text", ""))
            bg = f'background:{b["fill"]};' if b.get("fill") else ""
            content = (
                f'<div style="{bg}padding:6px;font-family:\'Noto Sans JP\',sans-serif;'
                f'font-size:13px;color:#111;white-space:pre-wrap;">{text_val}</div>'
            )

        html += f'<div style="{style}">{content}</div>'

    html += "</div></div>"
    return html

def parse_ai_output(text: str):
    """AI出力を6項目に分割"""
    sections = {
        "目標": "",
        "現状": "",
        "ビジネス課題": "",
        "調査目的": "",
        "問い": "",
        "仮説": "",
    }
    for key in sections.keys():
        pattern = rf"【{key}】(.*?)(?=【|$)"
        m = re.search(pattern, text, re.DOTALL)
        if m:
            sections[key] = m.group(1).strip()
    return sections



from pptx.enum.shapes import MSO_SHAPE_TYPE  # 既にインポート済みならこの行は重複していてもOK

from pptx.enum.shapes import MSO_SHAPE_TYPE  # 既にインポート済みならこの行は重複していてもOK
from pptx.dml.color import RGBColor         # ← これも上にあれば重複OK

def set_text_to_named_shape(slide, shape_name: str, text: str) -> bool:
    """
    スライド内の図形（グループ内も含む）から name=shape_name を探し、
    テキストを書き込む。見つかれば True、見つからなければ False を返す。
    - オートシェイプ／プレースホルダー：.text に書き込む
    - テーブル：全セルに同じテキストを書き込む（暫定）
    - 書き込んだテキストの文字色は黒（RGB 0,0,0）に設定する
    """

    def _set_font_black_textframe(text_frame):
        """text_frame 内の全 run のフォント色を黒にする"""
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font is not None:
                        run.font.color.rgb = RGBColor(0, 0, 0)
        except Exception:
            # フォーマット構造が想定外でも落ちないようにする
            pass

    def _search(shapes):
        for shp in shapes:
            # グループ内なら再帰
            if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                if _search(shp.shapes):
                    return True

            # 名前一致チェック
            if shp.name == shape_name:
                # テキスト枠があるタイプ
                if getattr(shp, "has_text_frame", False):
                    shp.text = text
                    _set_font_black_textframe(shp.text_frame)
                    return True

                # テーブルの場合
                if shp.shape_type == MSO_SHAPE_TYPE.TABLE:
                    try:
                        for row in shp.table.rows:
                            for cell in row.cells:
                                cell.text = text
                                _set_font_black_textframe(cell.text_frame)
                        return True
                    except Exception:
                        pass

        return False

    return _search(slide.shapes)


def parse_ai_output(text: str):
    """AI出力を6項目に分割"""
    sections = {
        "目標": "",
        "現状": "",
        "ビジネス課題": "",
        "調査目的": "",
        "問い": "",
        "仮説": "",
    }
    for key in sections.keys():
        pattern = rf"【{key}】(.*?)(?=【|$)"
        m = re.search(pattern, text, re.DOTALL)
        if m:
            sections[key] = m.group(1).strip()
    return sections


# ★ 調査仕様の項目（ラベルと session_state のキー）
SPEC_ITEMS = [
    ("調査手法", "spec_method"),
    ("抽出方法", "spec_sampling"),
    ("調査地域", "spec_region"),
    ("対象者条件", "spec_target"),
    ("サンプルサイズ", "spec_sample_size"),
    ("調査ボリューム", "spec_volume"),
    ("提示物", "spec_stimulus"),
    ("集計・分析仕様", "spec_analysis"),
    ("自由回答データの処理", "spec_openend"),
    ("業務範囲", "spec_scope"),
    ("納品物", "spec_deliverables"),
    ("インスペクションの方法", "spec_inspection"),
    ("謝礼の種類", "spec_incentive"),
    ("備考", "spec_notes"),
]

# ★ 調査仕様スライド（スライド6）の shape 名との対応
SPEC_LABEL_TO_SHAPE = {
    "調査手法": "Edit_SYUHO",
    "抽出方法": "Edit_Sampling",
    "調査地域": "Edit_Area",
    "対象者条件": "Edit_Joken",
    "サンプルサイズ": "Edit_Samplesize",
    "調査ボリューム": "Edit_Qvolume",
    "提示物": "Edit_review",
    "集計・分析仕様": "Edit_Analitics",
    "自由回答データの処理": "Edit_OAcdg",
    "業務範囲": "Edit_Hani",
    "納品物": "Edit_Nohin",
    "インスペクションの方法": "Edit_Inspection",
    "謝礼の種類": "Edit_Syarei",
    "備考": "Edit_Biko",
}

from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt

def apply_text_format(shape, font_name="Arial", font_size=12, color=RGBColor(0, 0, 0)):
    """
    shape.text_frame の paragraph/run に書式を統一的に適用する
    """
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT  # 左寄せ
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = color

import re

def parse_subquestions(ai_text: str):
    """
    『問いの分解』モードのAI出力をパースして、
    [
      {"subq": "...", "axis": "...", "metric": "..."},
      ...
    ]
    のリストに変換する
    """
    if not ai_text:
        return []

    # 「- サブクエスチョン...」でブロックごとに分割
    blocks = re.split(r"\n(?=-\s*サブクエスチョン)", ai_text.strip())
    results = []

    for blk in blocks:
        # サブクエスチョン本体
        m_q = re.search(r"-\s*サブクエスチョン[0-9０-９]*[:：]\s*(.+)", blk)
        if not m_q:
            continue

        # 分析軸
        m_axis = re.search(r"分析軸[:：]\s*(.+)", blk)
        # 評価項目
        m_metric = re.search(r"評価項目[:：]\s*(.+)", blk)

        results.append(
            {
                "subq": m_q.group(1).strip(),
                "axis": m_axis.group(1).strip() if m_axis else "",
                "metric": m_metric.group(1).strip() if m_metric else "",
            }
        )

    return results



from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_text_style(shape):
    """
    指定した図形内テキストの書式を統一するヘルパー
    - フォント：Arial
    - サイズ：12pt
    - 色：黒
    - 配置：左揃え
    """
    if not getattr(shape, "has_text_frame", False):
        return

    try:
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                if run.font is None:
                    continue
                run.font.name = "Arial"
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(0, 0, 0)
    except Exception:
        # 万一フォーマット構造が想定外でも、ここでは落とさない
        pass


# =========================
# レイアウト構成
# =========================
left, center, right = st.columns([1, 3, 1], gap="large")

# =========================
# 左ペイン
# =========================
with left:
    st.subheader("オリエン内容の整理")

    # ★追加：オリエン内容の整理ボタン
    if st.button("オリエン内容の整理", use_container_width=True):
        st.session_state["selected_mode"] = "オリエン内容の整理"
        st.session_state["message_center"] = ""
        st.session_state["message_right"] = ""
        st.rerun()

    st.subheader("市場/ブランドの整理")

    # --- 既存機能（ナビゲーション） ---
    if st.button("ブランド診断", use_container_width=True):
        st.session_state["selected_mode"] = "brand_diagnosis"
        st.session_state["message_center"] = ""
        st.session_state["message_right"] = ""
        st.rerun()

    st.divider()
    st.subheader("企画書構成")

    slide_names = [
        "表紙", "キックオフノート", "問いの分解","分析アプローチ",
        "対象者条件を検討", "調査項目案", "調査仕様案",
        "スケジュール案", "概算見積", "パワーポイントを出力"
    ]
    for idx, name in enumerate(slide_names):
        if st.button(name, use_container_width=True):
            st.session_state["selected_mode"] = name
            st.session_state["slide_index"] = idx  # スライド番号を保存
            st.session_state["message_center"] = ""
            st.session_state["message_right"] = ""
            st.rerun()

    st.divider()
    st.subheader("データ読み込み")

    # --- オリエン資料アップロード（トーン統一） ---
    uploaded_files = st.file_uploader(
        "オリエン資料をアップロードしてください（PDF / PPTX / TXT / ZIP）",
        type=["pdf", "pptx", "txt", "zip"],
        accept_multiple_files=True,
    )

    if uploaded_files:
        tempdir = tempfile.mkdtemp()
        texts = []
        for file in uploaded_files:
            path = os.path.join(tempdir, file.name)
            with open(path, "wb") as f:
                f.write(file.read())
            if path.endswith(".pdf"):
                texts.append(read_pdf(path))
            elif path.endswith(".pptx"):
                texts.append(read_pptx_text(path))
            elif path.endswith(".txt"):
                texts.append(read_txt(path))
            elif path.endswith(".zip"):
                with zipfile.ZipFile(path, "r") as z:
                    z.extractall(tempdir)
                for root, _, files in os.walk(tempdir):
                    for fn in files:
                        fp = os.path.join(root, fn)
                        if fn.endswith(".pdf"):
                            texts.append(read_pdf(fp))
                        elif fn.endswith(".pptx"):
                            texts.append(read_pptx_text(fp))
                        elif fn.endswith(".txt"):
                            texts.append(read_txt(fp))

        st.session_state["uploaded_docs"] = texts
        # 成功メッセージ（確定ではなく“共有・開始”のトーン）
        st.success(f"資料を共有しました。ここから一緒に読み解いていきましょう。（{len(uploaded_files)}件）")


    st.divider()

    # --- PPTテンプレートアップロード（トーン統一）---
    uploaded_pptx = st.file_uploader(
        "企画書テンプレートをアップロードしてください（PPTX）",
        type=["pptx"],
        key="pptx_upload"
    )

    # 初回アップロード時のみ pptx_path をセットする
    if uploaded_pptx and "template_loaded" not in st.session_state:
        cache_dir = Path(tempfile.gettempdir()) / "pptx_cache"
        cache_dir.mkdir(exist_ok=True)

        target = cache_dir / uploaded_pptx.name
        with open(target, "wb") as f:
            f.write(uploaded_pptx.getbuffer())

        st.session_state["pptx_path"] = str(target)
        st.session_state["template_loaded"] = True
        st.success(f"{uploaded_pptx.name} を読み込みました。")




# =========================
# 中央ペイン
# =========================
with center:
    mode = st.session_state.get("selected_mode")
    pptx_path = st.session_state.get("pptx_path")


    # =========================
    # 中央ペイン
    # === オリエン内容の整理 ===
    if mode == "オリエン内容の整理":
        st.markdown("## オリエン内容の整理")

        # 右ペインで生成したオリエン内容の下書き全文を表示
        # ★ウィジェットの key と、保存用の key を分ける
        if st.session_state.get("orien_outline_text"):
            # ★ 初回だけ、保存用テキスト → 編集用テキスト にコピーしておく
            if "orien_outline_editor" not in st.session_state:
                st.session_state["orien_outline_editor"] = st.session_state["orien_outline_text"]

            # ★ value を指定しない：Session State["orien_outline_editor"] が自動で入る
            st.text_area(
                "オリエン内容の下書き（編集可）",
                height=1500,
                key="orien_outline_editor",
            )

            # ★ 編集結果を保存用キーに反映
            st.session_state["orien_outline_text"] = st.session_state["orien_outline_editor"]

        else:
            st.info("右ペインの「下書き開始」ボタンを押すと、オリエン資料から抽出した下書きがここに表示されます。")

        # ❸ テキストをファイルに保存
        if st.button("この内容をテキストファイルとして保存", use_container_width=True):
            from pathlib import Path
            from datetime import datetime

            save_dir = get_session_dir()
            save_dir.mkdir(parents=True, exist_ok=True)

            filename = save_dir / f"orien_outline_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
            with open(filename, "w", encoding="utf-8") as f:
                f.write(st.session_state.get("orien_outline_text", ""))

            st.success(f"オリエン内容を保存しました：{filename}")


    # =========================
    # 中央ペイン
    # ====== ブランド診断モード ======
    elif st.session_state.get("selected_mode") == "brand_diagnosis":
        st.markdown("## ブランド診断")
        st.caption("右ペインでカテゴリーとブランドを推測・設定し、検索結果をここに表示します。")

        # --- 検索結果表示（AI出力がある場合） ---
        if "df_category_structure" in st.session_state:
            st.markdown("カテゴリーとブランドについて考察したこと")
            st.data_editor(
                st.session_state["df_category_structure"],
                hide_index=True,
                num_rows="fixed",
                use_container_width=True,
                key="editor_cat_struct_diag",
            )

        if "df_behavior_traits" in st.session_state:
            st.markdown("カテゴリーの消費行動特性として考察したこと")
            st.data_editor(
                st.session_state["df_behavior_traits"],
                hide_index=True,
                num_rows="fixed",
                use_container_width=True,
                key="editor_beh_traits_diag",
            )

        if (
            "df_category_structure" not in st.session_state
            and "df_behavior_traits" not in st.session_state
        ):
            st.info("右ペインで『カテゴリー・ブランドについて検索』を実行してください。")

        # --- マーケティングファネル表示 ---
        if "funnel_text" in st.session_state and st.session_state["funnel_text"]:
            st.markdown("### 🔄 マーケティングファネルとトリガー／障壁")
            st.markdown(f"```text\n{st.session_state['funnel_text']}\n```", unsafe_allow_html=True)
        else:
            st.info("右ペインで『カテゴリー・ブランドについて検索』を実行すると、ファネル情報がここに表示されます。")

        
    # =========================
    # 中央ペイン
    # === 表紙 ===
    elif mode == "表紙":
        st.markdown("## 表紙プレビュー（顧客名・調査名を自動推測）")

        pptx_path = st.session_state.get("pptx_path")

        # 🧠 AIで顧客名・調査名を自動推測
        ori_texts = "\n".join(st.session_state.get("uploaded_docs", []))
        if ori_texts and (
            not st.session_state.get("ai_client_name")
            or not st.session_state.get("ai_project_title")
        ):
            with st.spinner("顧客名と調査名を推測中..."):
                prompt = f"""
    あなたは市場調査の専門家です。
    以下のオリエン資料から、顧客企業名と調査タイトルを抽出・推定してください。

    【出力形式】
    顧客名：
    調査名：

    資料内容：
    {ori_texts[:4000]}
    """
                try:
                    response = client.chat.completions.create(
                        model=DEPLOYMENT,
                        messages=[
                            {"role": "system", "content": "あなたは市場調査の専門家です。"},
                            {"role": "user", "content": prompt},
                        ],
                        temperature=0.5,
                        max_tokens=200,
                    )
                    ai_result = response.choices[0].message.content
                    import re

                    client_match = re.search(r"顧客名[:：]\s*(.*)", ai_result)
                    title_match = re.search(r"調査名[:：]\s*(.*)", ai_result)

                    st.session_state["ai_client_name"] = (
                        client_match.group(1).strip() if client_match else ""
                    )
                    st.session_state["ai_project_title"] = (
                        title_match.group(1).strip() if title_match else ""
                    )

                    st.toast("顧客名・調査名を推測しました。", icon="🤖")
                except Exception as e:
                    st.error(f"AI呼び出しエラー: {e}")


        # 🖼 PowerPointプレビュー表示
        if pptx_path:
            from pptx import Presentation

            try:
                prs = Presentation(pptx_path)
                slide_index = 0

                if slide_index < len(prs.slides):
                    model = extract_slide_model(prs, slide_index=slide_index)
                    edited_texts = st.session_state.get("edited_texts", {})

                    # 入力済み or AI推測結果を反映
                    for key, ai_key in [
                        ("Edit_client", "ai_client_name"),
                        ("Edit_title", "ai_project_title"),
                    ]:
                        val = st.session_state.get(key)
                        if not val and st.session_state.get(ai_key):
                            st.session_state[key] = st.session_state[ai_key]
                            edited_texts[key] = st.session_state[ai_key]

                    if st.session_state.get("Edit_date"):
                        edited_texts["Edit_date"] = st.session_state["Edit_date"]

                    st.session_state["edited_texts"] = edited_texts

                    html = render_slide_html(model, edited_texts)
                    st.components.v1.html(html, height=520, scrolling=False)
                    st.caption("スライド1（表紙）のプレビューを表示中")
                else:
                    st.warning("スライド1が見つかりません。")
            except Exception as e:
                st.error(f"PPTプレビュー生成中にエラーが発生しました: {e}")
        else:
            st.info("PPTテンプレートをアップロードしてください。")

        st.markdown("---")


        # 📝 入力フォーム
        from datetime import datetime

        st.session_state["Edit_client"] = st.text_input(
            "顧客名",
            value=st.session_state.get("Edit_client", st.session_state.get("ai_client_name", "")),
            placeholder="例：株式会社〇〇",
        )

        st.session_state["Edit_title"] = st.text_input(
            "調査名",
            value=st.session_state.get("Edit_title", st.session_state.get("ai_project_title", "")),
            placeholder="例：〇〇市場における消費者意識調査",
        )

        st.session_state["Edit_date"] = st.text_input(
            "日付（YYYY年MM月DD日）",
            value=st.session_state.get(
                "Edit_date", datetime.now().strftime("%Y年%m月%d日")
            ),
        )

        st.markdown("---")


        # PowerPoint反映ボタン
        if st.button("📤 スライド1（表紙）にこの内容を反映", use_container_width=True):
            try:
                from pathlib import Path
                from datetime import datetime
                prs = Presentation(st.session_state.pptx_path)
                slide_index = 0
                if slide_index < len(prs.slides):
                    slide = prs.slides[slide_index]
                    mapping = {
                        "Edit_client": st.session_state["Edit_client"],
                        "Edit_title": st.session_state["Edit_title"],
                        "Edit_date": st.session_state["Edit_date"],
                    }

                    for shape_name, val in mapping.items():
                        shape = next(
                            (s for s in slide.shapes if s.name == shape_name), None
                        )
                        if shape:
                            shape.text = val
                            st.session_state.edited_texts[shape_name] = val

                    SLIDES_DIR = get_session_dir()
                    out_path = SLIDES_DIR / f"slide1_cover_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                    prs.save(out_path)
                    st.session_state.pptx_path = out_path
                    st.success("スライド1（表紙）に反映しました！")
                    st.rerun()
                else:
                    st.error("スライド1がテンプレートに存在しません。")
            except Exception as e:
                st.error(f"PowerPoint反映処理中にエラーが発生しました: {e}")

        
    # =========================
    # 中央ペイン
    # === キックオフノート ===
    elif mode == "キックオフノート":
        st.markdown("---")
        st.markdown("## キックオフノート")

        pptx_path = st.session_state.get("pptx_path")

        # ===============================
        # 🖼 スライド2のPPTプレビュー
        # ===============================
        if pptx_path:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            slide_index = 1
            if slide_index < len(prs.slides):
                model = extract_slide_model(prs, slide_index=slide_index)
                html = render_slide_html(model, st.session_state.edited_texts)
                st.components.v1.html(html, height=520, scrolling=False)
                st.caption("📊 現在のスライド2（キックオフノート）プレビュー")
            else:
                st.warning("スライド2がテンプレートに存在しません。")
        else:
            st.info("PPTテンプレートをアップロードしてください。")

        st.markdown("---")

        # ===============================
        # 🧾 AI出力 or 手入力フォーム
        # ===============================
        st.session_state.ai_目標 = st.text_area("① 目標（to be）", st.session_state.get("ai_目標", ""), height=100)
        st.session_state.ai_現状 = st.text_area("② 現状（as is）", st.session_state.get("ai_現状", ""), height=100)
        st.session_state.ai_ビジネス課題 = st.text_area("③ ビジネス課題", st.session_state.get("ai_ビジネス課題", ""), height=100)
        st.session_state.ai_調査目的 = st.text_area("④ 調査目的", st.session_state.get("ai_調査目的", ""), height=100)
        st.session_state.ai_問い = st.text_area("⑤ 問い", st.session_state.get("ai_問い", ""), height=100)
        st.session_state.ai_仮説 = st.text_area("⑥ 仮説", st.session_state.get("ai_仮説", ""), height=100)

        st.divider()

        # ===============================
        # 📤 PowerPoint反映ボタン
        # ===============================
        if st.button("📤 スライド2に反映（①〜⑥）", use_container_width=True):
            if pptx_path:
                try:
                    prs = Presentation(pptx_path)
                    slide_index = 1
                    if slide_index < len(prs.slides):
                        slide = prs.slides[slide_index]

                        mapping = {
                            "EDIT_TO_BE": st.session_state.ai_目標,
                            "EDIT_AS_IS": st.session_state.ai_現状,
                            "EDIT_PROBLEM": st.session_state.ai_ビジネス課題,
                            "EDIT_PURPOSE": st.session_state.ai_調査目的,
                            "EDIT_QUESTION": st.session_state.ai_問い,
                            "EDIT_HYPOTHESIS": st.session_state.ai_仮説,
                        }

                        for name, text in mapping.items():
                            shp = next((s for s in slide.shapes if s.name == name), None)
                            if shp:
                                shp.text = text
                                apply_text_format(shp)  # ← ← 書式統一を適用！
                                st.session_state.edited_texts[name] = text

                        # 保存（別名保存でキャッシュ衝突回避）
                        from datetime import datetime
                        SLIDES_DIR = get_session_dir()
                        out_path = SLIDES_DIR / f"kickoff_slide2_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                        prs.save(out_path)
                        st.session_state["pptx_path"] = str(out_path)

                        st.success("スライド2（キックオフノート）に反映しました！")
                        st.rerun()
                    else:
                        st.error("スライド2がテンプレートに存在しません。")
                except Exception as e:
                    st.error(f"PowerPoint反映中にエラーが発生しました: {e}")
            else:
                st.warning("PPTテンプレートをアップロードしてください。")

        

    # =========================
    # 中央ペイン
    # === 問いの分解（問いの分解ビュー）===
    elif mode == "問いの分解":
        st.markdown("## 問いの分解")

        pptx_path = st.session_state.get("pptx_path")

        # ---- PPTプレビュー表示（HTMLプレビュー版：スライド3）----
        if pptx_path:
            try:
                from pptx import Presentation

                prs = Presentation(pptx_path)
                slide_index = 2  # スライド3（0始まり）

                if slide_index < len(prs.slides):
                    model = extract_slide_model(prs, slide_index=slide_index)
                    html = render_slide_html(model, st.session_state.get("edited_texts", {}))
                    st.components.v1.html(html, height=520, scrolling=False)
                    st.caption("📊 現在のスライド3（問いの分解）プレビュー")
                else:
                    st.warning("スライド3がテンプレートに存在しません。")
            except Exception as e:
                st.error(f"PPTプレビュー生成中にエラーが発生しました: {e}")
        else:
            st.info("PPTテンプレートをアップロードしてください。")

        st.markdown("---")

        # =========================================
        # ① 構造ビュー：目的 → メインクエスチョン → サブクエスチョン
        # =========================================
        import re

        main_question_text = st.session_state.get("ai_問い", "")
        purpose = st.session_state.get("ai_調査目的", "")
        subq_list = st.session_state.get("subq_list", [])

        st.markdown("### 構造ビュー：目的 → メインクエスチョン → サブクエスチョン")

        def split_main_questions(text: str):
            """
            キックオフノート⑤『問い』のテキストから
            「1. …」「2) …」「Q1: …」のような行を検出して
            メインクエスチョンのリストに分割する。
            """
            if not text:
                return []

            lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
            questions = []
            buf = ""

            for line in lines:
                # 先頭が「1.」「2)」「Q1:」などの行を新しい問いの開始とみなす
                m = re.match(r'^(?:\d+[\.\)]|Q\d+[:：])\s*(.+)', line)
                if m:
                    # 前のバッファを1問として確定
                    if buf:
                        questions.append(buf.strip())
                    buf = m.group(1)
                else:
                    # 前行の続き
                    if buf:
                        buf += " " + line
                    else:
                        buf = line

            if buf:
                questions.append(buf.strip())

            # 番号付きで取れなかった場合は全文を1問として扱う
            if not questions and text.strip():
                questions = [text.strip()]

            return questions

        # メインクエスチョン群を抽出
        main_questions = split_main_questions(main_question_text)

        tree_lines = []

        if not main_questions and not subq_list:
            st.info("右ペイン『問いの分解』でサブクエスチョンを生成すると、ここに構造ビューが表示されます。")
        else:
            # -------------------
            # 目的
            # -------------------
            tree_lines.append("目的（キックオフノート）")
            if purpose:
                tree_lines.append(f"  ┗ {purpose}")
            else:
                tree_lines.append("  ┗ （未設定）")

            tree_lines.append("")

            # -------------------
            # メインクエスチョン → サブクエスチョン
            # -------------------
            tree_lines.append("問い（メインクエスチョン）")

            if not main_questions:
                # メインQが1つも抽出できなかった場合：従来どおり1ブロックとして表示
                if main_question_text:
                    tree_lines.append(f"  ┗ {main_question_text}")
                    if subq_list:
                        for i, sq in enumerate(subq_list, 1):
                            tree_lines.append(f"       ┗ SQ{i}: {sq.get('subq', '')}")
                    else:
                        tree_lines.append("       ┗ （まだサブクエスチョンが生成されていません）")
                else:
                    tree_lines.append("  ┗ （未設定）")
            else:
                # メインQが複数ある場合：サブQをラウンドロビンで割り当て
                grouped = {mq: [] for mq in main_questions}

                if subq_list:
                    for idx, sq in enumerate(subq_list):
                        mq = sq.get("main_question")
                        if mq and mq in grouped:
                            # すでに main_question が付いていればそれを優先
                            grouped[mq].append(sq)
                        else:
                            # 付いていなければ順番に割り当て
                            target_mq = main_questions[idx % len(main_questions)]
                            grouped[target_mq].append(sq)

                # メインQごとにツリー表示
                for mq in main_questions:
                    tree_lines.append(f"  ┗ {mq}")
                    sq_items = grouped.get(mq, [])
                    if not sq_items:
                        tree_lines.append("       ┗ （まだサブクエスチョンが紐付いていません）")
                    else:
                        for i, sq in enumerate(sq_items, 1):
                            tree_lines.append(f"       ┗ SQ{i}: {sq.get('subq', '')}")
                    tree_lines.append("")

            st.code("\n".join(tree_lines), language="text")

        # 構造ビューのテキスト（PPT反映用）
        tree_text = "\n".join(tree_lines) if tree_lines else ""

        st.markdown("---")

        # =========================================
        # ② PPT出力：スライド3の EDIT1_subQ に反映
        #     → 構造ビューのテキスト内容を反映
        # =========================================
        if st.button("📤 この内容をスライド3（EDIT1_subQ）に反映", use_container_width=True):
            pptx_path = st.session_state.get("pptx_path")
            if not tree_text.strip():
                st.warning("構造ビューの内容が空です。先に『問い』やサブクエスチョンを設定してください。")
            elif pptx_path:
                try:
                    from pptx import Presentation
                    from pathlib import Path
                    from datetime import datetime

                    prs = Presentation(pptx_path)
                    slide_index = 2  # スライド3（0始まり）

                    if slide_index < len(prs.slides):
                        slide = prs.slides[slide_index]

                        # 構造ビューのテキストをそのまま反映
                        text_to_apply = tree_text

                        shp = next((s for s in slide.shapes if s.name == "EDIT1_subQ"), None)
                        if shp and getattr(shp, "has_text_frame", False):
                            shp.text = text_to_apply

                            # 共通フォーマット適用（黒・12pt・左寄せ・Arial）
                            try:
                                apply_text_style(shp)
                            except Exception:
                                # ヘルパー側でエラーになってもアプリが落ちないように
                                pass

                            # プレビュー用キャッシュ
                            st.session_state.edited_texts["EDIT1_subQ"] = text_to_apply
                            st.session_state.edited_texts["EDIT1_QUESTION_FACTORS"] = text_to_apply

                            # 保存（別名保存でキャッシュ衝突回避）
                            SLIDES_DIR = get_session_dir()
                            SLIDES_DIR.mkdir(parents=True, exist_ok=True)
                            out_path = SLIDES_DIR / f"question_factors_slide3_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                            prs.save(out_path)
                            st.session_state["pptx_path"] = str(out_path)

                            st.success("スライド3（問いの分解）に構造ビューの内容を反映しました！（EDIT1_subQ・書式統一）")
                            st.rerun()
                        else:
                            st.error("スライド3に名前が『EDIT1_subQ』のテキスト図形が見つかりませんでした。")
                    else:
                        st.error("スライド3がテンプレートに存在しません。")

                except Exception as e:
                    st.error(f"PowerPoint反映中にエラーが発生しました: {e}")
            else:
                st.warning("PPTテンプレートをアップロードしてください。")


    # =========================
    # 中央ペイン
    # === 分析アプローチ ===
    elif mode == "分析アプローチ":
        st.markdown("## 分析アプローチ")

        pptx_path = st.session_state.get("pptx_path")

        # ===============================
        # ① PPTプレビュー（上部に表示）
        #    スライド4〜12をタブで切替表示
        # ===============================
        if pptx_path:
            try:
                from pptx import Presentation

                prs = Presentation(pptx_path)
                edited_texts = st.session_state.get("edited_texts", {})

                # SQ1〜SQ9としてタブ表示（スライド4〜12）
                slide_numbers = list(range(1, 10))  # SQ1〜SQ9
                tab_labels = [f"SQ{n}" for n in slide_numbers]
                tabs = st.tabs(tab_labels)

                for idx, tab in enumerate(tabs):
                    slide_index = 3 + idx  # 0始まり → 3=スライド4
                    with tab:
                        if slide_index < len(prs.slides):
                            model = extract_slide_model(prs, slide_index=slide_index)
                            html = render_slide_html(model, edited_texts)
                            st.components.v1.html(html, height=520, scrolling=False)
                            st.caption(f"📊 現在のスライド{slide_index+1}（分析アプローチ）プレビュー")
                        else:
                            st.warning(f"スライド{slide_index+1}（分析アプローチ用）がテンプレートに存在しません。")
            except Exception as e:
                st.error(f"PPTプレビュー生成中にエラーが発生しました: {e}")
        else:
            st.info("PPTテンプレートをアップロードしてください。")

        st.markdown("---")

        # ===============================
        # ② サブクエスチョン別 5項目セット表示
        #    （サブQ / 分析軸 / 評価項目 / アプローチ / 仮説）
        # ===============================
        analysis_blocks = st.session_state.get("analysis_blocks", None)
        subq_list = st.session_state.get("subq_list", [])

        if not subq_list:
            st.info("まだサブクエスチョンが構造化されていません。先に『問いの分解』モードでサブクエスチョンを生成してください。")
        elif not analysis_blocks:
            st.info("右ペインの『🪄 下書きを作成』ボタンを押すと、サブクエスチョンごとの分析アプローチ案がここに表示されます。")
        else:
            st.markdown("### サブクエスチョン別 分析アプローチ案（5項目セット）")
            st.caption("1つのサブクエスチョンにつき、サブQ・分析軸・評価項目・アプローチ・仮説を個別に編集できます。")

            # 旧ロジックの痕跡をクリア（任意）
            if "analysis_block_texts" in st.session_state:
                del st.session_state["analysis_block_texts"]

            # サブQごとに 5項目の入力欄を表示
            for i, blk in enumerate(analysis_blocks, 1):
                st.markdown(f"#### サブクエスチョン {i}")

                # セッションに保持するキー
                subq_key = f"analysis_subq_{i}"
                axis_key = f"analysis_axis_{i}"
                metric_key = f"analysis_metric_{i}"
                approach_key = f"analysis_approach_{i}"
                hypo_key = f"analysis_hypothesis_{i}"

                # ウィジェット用キー（表示側）
                subq_widget_key = f"{subq_key}_input"
                axis_widget_key = f"{axis_key}_input"
                metric_widget_key = f"{metric_key}_input"
                approach_widget_key = f"{approach_key}_input"
                hypo_widget_key = f"{hypo_key}_input"

                # 初期値：session_state にあればそれを優先、なければ AI の blk から
                subq_val = st.text_area(
                    "サブクエスチョン",
                    value=st.session_state.get(subq_key, blk.get("subq", "")),
                    height=80,
                    key=subq_widget_key,
                )
                axis_val = st.text_area(
                    "分析軸（セグメント）",
                    value=st.session_state.get(axis_key, blk.get("axis", "")),
                    height=60,
                    key=axis_widget_key,
                )
                metric_val = st.text_area(
                    "評価項目",
                    value=st.session_state.get(metric_key, blk.get("metric", "")),
                    height=60,
                    key=metric_widget_key,
                )
                approach_val = st.text_area(
                    "主な分析アプローチ",
                    value=st.session_state.get(approach_key, blk.get("approach", "")),
                    height=80,
                    key=approach_widget_key,
                )
                hypo_val = st.text_area(
                    "検証する仮説",
                    value=st.session_state.get(hypo_key, blk.get("hypothesis", "")),
                    height=80,
                    key=hypo_widget_key,
                )

                # 入力値を session_state に確定保存（モード切替しても残るように）
                st.session_state[subq_key] = subq_val
                st.session_state[axis_key] = axis_val
                st.session_state[metric_key] = metric_val
                st.session_state[approach_key] = approach_val
                st.session_state[hypo_key] = hypo_val

                st.markdown("---")

            # edited_texts を PPTビューア用に更新
            edited_texts = st.session_state.get("edited_texts", {})

            # analysis_blocks 自体も上書き（必要なら）
            for i, blk in enumerate(analysis_blocks, 1):
                subq_key = f"analysis_subq_{i}"
                axis_key = f"analysis_axis_{i}"
                metric_key = f"analysis_metric_{i}"
                approach_key = f"analysis_approach_{i}"
                hypo_key = f"analysis_hypothesis_{i}"

                subq = st.session_state.get(subq_key, "")
                axis = st.session_state.get(axis_key, "")
                metric = st.session_state.get(metric_key, "")
                approach = st.session_state.get(approach_key, "")
                hypothesis = st.session_state.get(hypo_key, "")

                blk["subq"] = subq
                blk["axis"] = axis
                blk["metric"] = metric
                blk["approach"] = approach
                blk["hypothesis"] = hypothesis

                # PPTビューアー用：Shape名ごとに格納
                idx = i  # サブQ番号
                edited_texts[f"EDIT1_subQ{idx}_1"] = subq
                edited_texts[f"EDIT1_subQ{idx}_2"] = axis
                edited_texts[f"EDIT1_subQ{idx}_3"] = metric
                edited_texts[f"EDIT1_subQ{idx}_4"] = approach
                edited_texts[f"EDIT1_subQ{idx}_5"] = hypothesis

            st.session_state["analysis_blocks"] = analysis_blocks
            st.session_state["edited_texts"] = edited_texts

            # ===============================
            # ③ PPT 本体に反映するボタン
            #    スライド4〜12にサブQごと・項目ごとに分割して反映
            # ===============================
            if st.button("📤 この内容をスライド4に反映（サブQごとにスライド分割）", use_container_width=True):
                pptx_path = st.session_state.get("pptx_path")
                if not pptx_path:
                    st.warning("PPTテンプレートを先にアップロードしてください。")
                else:
                    try:
                        from pptx import Presentation
                        from pathlib import Path
                        from datetime import datetime

                        prs = Presentation(pptx_path)

                        max_slides = 9  # スライド4〜12 → 最大9サブクエスチョン
                        total_blocks = len(analysis_blocks)

                        if total_blocks > max_slides:
                            st.warning(
                                f"サブクエスチョンが {total_blocks} 個ありますが、"
                                f"スライドは最大 {max_slides} 枚までの対応です。先頭 {max_slides} 件のみを反映します。"
                            )

                        applied_count = 0

                        for i in range(1, min(total_blocks, max_slides) + 1):
                            slide_index = 3 + (i - 1)  # スライド4〜12 → index=3〜11
                            if slide_index >= len(prs.slides):
                                st.warning(
                                    f"テンプレート内のスライド数が不足しているため、"
                                    f"サブクエスチョン{i}以降は反映できませんでした。"
                                )
                                break

                            slide = prs.slides[slide_index]

                            subq = st.session_state.get(f"analysis_subq_{i}", "")
                            axis = st.session_state.get(f"analysis_axis_{i}", "")
                            metric = st.session_state.get(f"analysis_metric_{i}", "")
                            approach = st.session_state.get(f"analysis_approach_{i}", "")
                            hypothesis = st.session_state.get(f"analysis_hypothesis_{i}", "")

                            values = [
                                (f"EDIT1_subQ{i}_1", subq),
                                (f"EDIT1_subQ{i}_2", axis),
                                (f"EDIT1_subQ{i}_3", metric),
                                (f"EDIT1_subQ{i}_4", approach),
                                (f"EDIT1_subQ{i}_5", hypothesis),
                            ]

                            for shape_name, text_val in values:
                                if not text_val:
                                    continue

                                ok = set_text_to_named_shape(slide, shape_name, text_val)

                                if ok:
                                    shp = next((s for s in slide.shapes if s.name == shape_name), None)
                                    if shp and getattr(shp, "has_text_frame", False):
                                        apply_text_format(shp)

                                    # プレビュー用キャッシュも更新
                                    st.session_state.edited_texts[shape_name] = text_val
                                    applied_count += 1
                                else:
                                    st.warning(
                                        f"スライド{slide_index+1}内に '{shape_name}' という名前のテキスト図形が見つかりませんでした。"
                                    )

                        if applied_count > 0:
                            SLIDES_DIR = get_session_dir()
                            SLIDES_DIR.mkdir(parents=True, exist_ok=True)
                            out_path = SLIDES_DIR / f"analysis_approach_slide4to12_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                            prs.save(out_path)

                            st.session_state["pptx_path"] = str(out_path)

                            st.success(
                                f"スライド4〜12（分析アプローチ）にサブクエスチョン別・項目別の内容を反映しました！（{applied_count}箇所）"
                            )
                            st.rerun()
                        else:
                            st.error("いずれのスライドにもテキストを反映できませんでした。Shape名やテンプレート構成を確認してください。")

                    except Exception as e:
                        st.error(f"PowerPoint反映中にエラーが発生しました: {e}")

                
        
    # =========================
    # 中央ペイン
    # === 対象者条件を検討 ===
    elif mode == "対象者条件を検討":
        st.markdown("## 対象者条件を検討")

        pptx_path = st.session_state.get("pptx_path")

        # ---- PPTプレビュー表示（HTMLプレビュー版：スライド4）----
        if pptx_path:
            try:
                from pptx import Presentation

                prs = Presentation(pptx_path)
                slide_index = 12  # スライド4（0始まり）

                if slide_index < len(prs.slides):
                    # ★ 問いの要因分解と同じ：extract → render
                    model = extract_slide_model(prs, slide_index=slide_index)
                    html = render_slide_html(model, st.session_state.get("edited_texts", {}))
                    st.components.v1.html(html, height=520, scrolling=False)
                    st.caption("📊 現在のスライド4（対象者条件を検討）プレビュー")
                else:
                    st.warning("スライド4がテンプレートに存在しません。")
            except Exception as e:
                st.error(f"PPTプレビュー生成中にエラーが発生しました: {e}")
        else:
            st.info("PPTテンプレートをアップロードしてください。")

        st.markdown("---")

        # ---- AI生成結果を表示 ＋ PPT出力 ----
        if st.session_state.get("ai_target_condition"):
            st.markdown("### 調査対象者条件案")
            st.text_area(
                "検討した対象者条件（編集可）",
                value=st.session_state["ai_target_condition"],
                height=300,
                key="target_condition_textarea"
            )

            st.markdown("---")

            if st.button("📤 この内容をスライド4（EDIT1_taisyosya）に反映", use_container_width=True):
                pptx_path = st.session_state.get("pptx_path")

                if pptx_path:
                    try:
                        from pptx import Presentation
                        from pathlib import Path
                        from datetime import datetime

                        prs = Presentation(pptx_path)
                        slide_index = 12  # スライド4（0始まり）

                        if slide_index < len(prs.slides):
                            slide = prs.slides[slide_index]

                            # テキストエリアの編集内容を優先
                            text_to_apply = st.session_state.get(
                                "target_condition_textarea",
                                st.session_state.get("ai_target_condition", "")
                            )

                            # EDIT1_taisyosya を探す
                            shp = next((s for s in slide.shapes if s.name == "EDIT1_taisyosya"), None)
                            if shp and getattr(shp, "has_text_frame", False):

                                # ★ テキストを反映
                                shp.text = text_to_apply

                                # ★ 統一書式を適用（Arial / 12pt / 左寄せ / 黒）
                                apply_text_format(shp)

                                # プレビュー用キャッシュ
                                st.session_state.edited_texts["EDIT1_taisyosya"] = text_to_apply
                                st.session_state.edited_texts["EDIT1_TARGET_CONDITION"] = text_to_apply

                                # 保存（別名保存）
                                SLIDES_DIR = get_session_dir()
                                SLIDES_DIR.mkdir(parents=True, exist_ok=True)
                                out_path = SLIDES_DIR / f"target_condition_slide4_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

                                prs.save(out_path)
                                st.session_state["pptx_path"] = str(out_path)

                                st.success("スライド4（対象者条件）に反映しました！（フォント・色・左寄せを統一）")
                                st.rerun()

                            else:
                                st.error("スライド4に『EDIT1_taisyosya』のテキスト図形が見つかりませんでした。")

                        else:
                            st.error("スライド4がテンプレートに存在しません。")

                    except Exception as e:
                        st.error(f"PowerPoint反映中にエラーが発生しました: {e}")

                else:
                    st.warning("PPTテンプレートをアップロードしてください。")

        else:
            st.info("右ペインの『🪄 下書きを作成』を押すと、対象者条件案がここに表示されます。")


    # =========================
    # 中央ペイン
    # === 調査項目案 ===
    elif mode == "調査項目案":
        st.markdown("## 調査項目案")

        pptx_path = st.session_state.get("pptx_path")

        # ---- PPTプレビュー表示（HTMLプレビュー版）----
        if pptx_path:
            try:
                from pptx import Presentation

                prs = Presentation(pptx_path)
                slide_index = 13  # スライド13（0始まり）

                if slide_index < len(prs.slides):
                    model = extract_slide_model(prs, slide_index=slide_index)
                    html = render_slide_html(model, st.session_state.get("edited_texts", {}))
                    st.components.v1.html(html, height=520, scrolling=False)
                    st.caption("📊 現在のスライド5（調査項目案）プレビュー")
                else:
                    st.warning("スライド5がテンプレートに存在しません。")
            except Exception as e:
                st.error(f"PPTプレビュー生成中にエラーが発生しました: {e}")
        else:
            st.info("PPTテンプレートをアップロードしてください。")

        st.markdown("---")

        # ---- AI生成結果を表示 ----
        if st.session_state.get("ai_survey_items"):
            st.markdown("### 🤖 AI生成：調査項目案（バリエーション別）")

            items = st.session_state["ai_survey_items"]

            for ver in ["10問", "20問", "30問", "40問"]:
                text_key = f"survey_items_{ver}"
                default_val = items.get(ver, "")

                if default_val:
                    # 🔹テキストエリア（編集可）
                    st.text_area(
                        f"📝 {ver}バージョン",
                        value=default_val,
                        height=500,
                        key=text_key,
                    )

                    # ===============================
                    # 📤 PowerPoint反映ボタン（EDIT1）
                    # ===============================
                    if st.button(
                        f"📤 この{ver}バージョンをスライド5に反映（EDIT1_Qimg）",
                        use_container_width=True,
                        key=f"apply_{ver}",
                    ):
                        pptx_path = st.session_state.get("pptx_path")
                        if pptx_path:
                            try:
                                from pptx import Presentation
                                from pathlib import Path
                                from datetime import datetime

                                prs = Presentation(pptx_path)
                                slide_index = 13  # スライド13（0始まり）

                                if slide_index < len(prs.slides):
                                    slide = prs.slides[slide_index]

                                    # 最新のテキストエリア内容を取得
                                    text_to_apply = st.session_state.get(text_key, default_val)

                                    # EDIT1_Qimg を直接探して text を代入
                                    shp = next((s for s in slide.shapes if s.name == "EDIT1_Qimg"), None)
                                    if shp and getattr(shp, "has_text_frame", False):

                                        # ★ テキストを反映
                                        shp.text = text_to_apply

                                        # ★ 統一書式を適用（Arial / 12pt / 黒 / 左寄せ）
                                        apply_text_format(shp)

                                        # プレビュー用にも保存
                                        st.session_state.edited_texts["EDIT1_Qimg"] = text_to_apply
                                        st.session_state.edited_texts["EDIT1_SURVEY_ITEMS"] = text_to_apply

                                        # 保存（別名保存でキャッシュ衝突回避）
                                        SLIDES_DIR = get_session_dir()
                                        SLIDES_DIR.mkdir(parents=True, exist_ok=True)

                                        out_path = SLIDES_DIR / f"surveyitems_slide5_{ver}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                                        prs.save(out_path)
                                        st.session_state["pptx_path"] = str(out_path)

                                        st.success(
                                            f"スライド5（調査項目案）に {ver} バージョンを反映しました！（フォント・サイズ・色を統一）"
                                        )
                                        st.rerun()

                                    else:
                                        st.error("スライド5に『EDIT1_Qimg』という名前のテキスト図形が見つかりませんでした。")

                                else:
                                    st.error("スライド5がテンプレートに存在しません。")

                            except Exception as e:
                                st.error(f"PowerPoint反映中にエラーが発生しました: {e}")
                        else:
                            st.warning("PPTテンプレートをアップロードしてください。")


    # =========================
    # 中央ペイン
    # === 調査仕様案 ===
    elif mode == "調査仕様案":
        st.markdown("## 調査仕様案")

        pptx_path = st.session_state.get("pptx_path")

        # ---- PPTプレビュー表示（HTMLプレビュー版：スライド6）----
        if pptx_path:
            try:
                from pptx import Presentation

                prs = Presentation(pptx_path)
                slide_index = 14  # スライド14（0始まり）

                if slide_index < len(prs.slides):
                    model = extract_slide_model(prs, slide_index=slide_index)
                    html = render_slide_html(model, st.session_state.get("edited_texts", {}))
                    st.components.v1.html(html, height=520, scrolling=False)
                    st.caption("📊 現在のスライド6（調査仕様案）プレビュー")
                else:
                    st.warning("スライド6がテンプレートに存在しません。")
            except Exception as e:
                st.error(f"PPTプレビュー生成中にエラーが発生しました: {e}")
        else:
            st.info("PPTテンプレートをアップロードしてください。")

        st.markdown("---")

        # ---- 調査仕様の下書き表示 ----
        has_spec = any(st.session_state.get(key) for _, key in SPEC_ITEMS)

        if has_spec:
            st.markdown("### 調査仕様の下書き（編集可）")

            # 仕様項目のテキストエリア
            for label, key in SPEC_ITEMS:
                st.text_area(
                    label,
                    height=80,
                    key=key,
                )

            st.markdown("---")

            # ===============================
            # 📤 調査仕様をPPT（スライド6）に反映
            # ===============================
            if st.button("📤 この調査仕様をスライド6に反映", use_container_width=True):
                pptx_path = st.session_state.get("pptx_path")
                if not pptx_path:
                    st.warning("PPTテンプレートを先にアップロードしてください。")
                else:
                    try:
                        from pptx import Presentation
                        from pathlib import Path
                        from datetime import datetime

                        prs = Presentation(pptx_path)
                        slide_index = 14  # スライド14（0始まり）

                        if slide_index < len(prs.slides):
                            slide = prs.slides[slide_index]

                            # ==========================
                            # 仕様項目 → PPT shape へ反映
                            # ==========================
                            for label, key in SPEC_ITEMS:
                                text_val = st.session_state.get(key, "")
                                shape_name = SPEC_LABEL_TO_SHAPE.get(label)

                                if shape_name and text_val is not None:

                                    # shape へ書き込む（set_text_to_named_shape: グループ対応）
                                    ok = set_text_to_named_shape(slide, shape_name, text_val)

                                    if ok:
                                        # shape を再取得して書式を適用
                                        shp = next((s for s in slide.shapes if s.name == shape_name), None)

                                        if shp and getattr(shp, "has_text_frame", False):

                                            # ★ 統一書式を適用（Arial / 12pt / 黒 / 左寄せ）
                                            apply_text_format(shp)

                                        # プレビュー用キャッシュ
                                        st.session_state.edited_texts[shape_name] = text_val

                            # 編集内容まとめ
                            st.session_state.edited_texts["SPEC_LAST_APPLIED"] = {
                                label: st.session_state.get(key, "")
                                for label, key in SPEC_ITEMS
                            }

                            # ==========================
                            # 保存（別名保存）
                            # ==========================
                            SLIDES_DIR = get_session_dir()
                            SLIDES_DIR.mkdir(parents=True, exist_ok=True)

                            out_path = SLIDES_DIR / f"spec_slide6_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                            prs.save(out_path)

                            st.session_state["pptx_path"] = str(out_path)
                            st.success("スライド6（調査仕様案）に調査仕様を反映しました！（フォント・色・左寄せを統一）")
                            st.rerun()

                        else:
                            st.error("スライド6がテンプレートに存在しません。")

                    except Exception as e:
                        st.error(f"PowerPoint反映中にエラーが発生しました: {e}")


    # =========================
    # 中央ペイン
    # === スケジュール案 ===
    elif mode == "スケジュール案":
        st.markdown("## スケジュール案")

        pptx_path = st.session_state.get("pptx_path")

        # ---- PPTプレビュー表示（スライド7：画像プレビュー）----
        if pptx_path:
            try:
                # 他モードと同じ方式：PPTX → 画像化して表示
                images = pptx_to_images(pptx_path)
                if len(images) > 15:  # スライド14は 0始まりで index=14
                    st.image(images[15], caption="スライド14：スケジュール案", use_container_width=True)
                else:
                    st.warning("スライド14がテンプレートに存在しません。")
            except Exception as e:
                st.error(f"PPTプレビュー生成中にエラーが発生しました: {e}")
        else:
            st.info("PPTテンプレートをアップロードしてください。")

        st.markdown("---")
        st.markdown("### 🗓 スケジュール案（工程・期間設定）")
        st.caption("『オリエン内容の整理』から抽出したマイルストンを起点に、工程ごとの日程を設定します。")

        # ====== ここからスケジュール計算ロジック ======
        from pptx import Presentation  # 重複インポートOK
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pandas.tseries.offsets import CustomBusinessDay
        from datetime import datetime
        import pandas as pd
        import pytz
        from io import BytesIO
        from pathlib import Path

        JST = pytz.timezone("Asia/Tokyo")

        # ------------------------------------------------
        # 祝日・年末年始を含む休暇設定
        # ------------------------------------------------
        import pandas as pd
        from datetime import datetime
        import jpholiday  # 事前に pip install jpholiday

        # ------------------------------------------------
        # 演算開始日から1年先までの祝日リストを生成
        # ------------------------------------------------
        def get_holiday_list(base_date=None, years=1):
            """
            base_date から years 年先までの期間について、
            日本の祝日 + 年末年始を pandas.Timestamp のリストで返す
            """
            if base_date is None:
                base_date = datetime.now().date()

            start = pd.Timestamp(base_date)
            end = start + pd.DateOffset(years=years)

            # 日本の祝日（jpholiday を使用）
            days = pd.date_range(start, end, freq="D")
            holidays = [d for d in days if jpholiday.is_holiday(d)]

            # 年末年始（任意で追加したい場合）
            # 例：その年と翌年の 12/29〜1/3 を毎回休暇扱いする
            for y in range(start.year, end.year + 1):
                newyear_span = pd.date_range(f"{y}-12-29", f"{y+1}-01-03", freq="D")
                holidays += list(newyear_span)

            # normalize して重複を除去
            holidays = sorted(list(set([pd.Timestamp(h).normalize() for h in holidays])))
            return holidays

        # ------------------------------------------------
        # 営業日スケジュール生成（順序固定）
        # ------------------------------------------------
        from pandas.tseries.offsets import CustomBusinessDay
        from datetime import datetime

        def compute_schedule_with_fixed(df_phase, base_date=None):
            # 基準日（None の場合は今日）
            base_dt = pd.to_datetime(base_date or datetime.now())

            # ← ここで基準日から1年先分の祝日リストを作る
            holidays = get_holiday_list(base_dt, years=1)

            biz_day = CustomBusinessDay(weekmask="Mon Tue Wed Thu Fri", holidays=holidays)

            df = df_phase.copy()
            df["fixed_date"] = pd.to_datetime(df["fixed_date"], errors="coerce")
            df["duration_days"] = df["duration_days"].fillna(1).astype(int)

            rows = []
            cur = base_dt
            last_end = cur

            for _, r in df.iterrows():
                dur = int(r["duration_days"])
                if pd.notna(r["fixed_date"]):
                    start = r["fixed_date"]
                    end = start + (dur - 1) * biz_day
                else:
                    start = max(last_end, cur)
                    end = start + (dur - 1) * biz_day

                rows.append({
                    "工程": r["name"],
                    "開始日": start,
                    "終了日": end,
                    "固定日": r["fixed_date"]
                })
                last_end = end
                cur = end + 1 * biz_day

            return pd.DataFrame(rows)

        # ------------------------------------------------
        # カレンダー表作成
        # ------------------------------------------------
        def make_calendar_table(schedule_df):
            youbi_jp = {"Mon": "月", "Tue": "火", "Wed": "水", "Thu": "木", "Fri": "金", "Sat": "土", "Sun": "日"}

            holidays = get_holiday_list()
            holidays_set = set(pd.to_datetime(holidays).strftime("%m/%d").tolist())

            start = schedule_df["開始日"].min()
            end = schedule_df["終了日"].max()
            days = pd.date_range(start, end, freq="D")

            cal_df = pd.DataFrame({
                "日付": days.strftime("%m/%d"),
                "曜日": days.strftime("%a").map(youbi_jp),
                "マイルストン": "",
                "非営業日": False
            })

            for _, r in schedule_df.iterrows():
                cal_df.loc[cal_df["日付"] == r["開始日"].strftime("%m/%d"), "マイルストン"] = r["工程"]

            cal_df.loc[
                (cal_df["曜日"].isin(["土", "日"])) | (cal_df["日付"].isin(holidays_set)),
                "非営業日"
            ] = True

            max_rows = 60
            if len(cal_df) > max_rows:
                cal_df = cal_df.iloc[:max_rows]

            return cal_df

        # ------------------------------------------------
        # （参考）単独スケジュールPPT作成用関数（今は使わなくてもOK）
        # ------------------------------------------------
        def calendar_to_pptx(df):
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
            p = txBox.text_frame.add_paragraph()
            p.font.bold = True
            p.font.size = Pt(14)
            p.text = "市場調査スケジュール表"

            rows, cols = df.shape[0] + 1, 3
            left, top, width, height = Inches(0.5), Inches(0.8), Inches(8.5), Inches(5.5)
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            headers = ["日付", "曜日", "マイルストン"]
            for j, col in enumerate(headers):
                cell = table.cell(0, j)
                cell.text = col
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.size = Pt(11)
                    p.alignment = PP_ALIGN.CENTER

            for i in range(df.shape[0]):
                for j, col in enumerate(headers):
                    val = df.at[i, col]
                    cell = table.cell(i + 1, j)
                    cell.text = str(val) if val else ""
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(11)
                    if df.at[i, "非営業日"]:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(220, 220, 220)
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            bio = BytesIO()
            prs.save(bio)
            bio.seek(0)
            return bio.read()

        # ------------------------------------------------
        # 🔑 スライド7にスケジュール表を反映する関数
        # ------------------------------------------------
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import pandas as pd
        import math

        def reflect_schedule_to_slide7(prs, calendar_df: pd.DataFrame):
            """
            スライド7にスケジュール表を3分割して挿入
            - calendar_df: 「日付」「曜日」「マイルストン」「非営業日」を含む DataFrame を想定
            - スライド上の Shape名 schedule1 / schedule2 / schedule3 の位置・サイズに表を配置
            - 非営業日(True)の行は薄いグレーでハイライト
            """
            slide_index = 15  # スライド7（0始まり）
            if slide_index >= len(prs.slides):
                st.error("スライド7がテンプレートに存在しません。")
                return prs

            slide = prs.slides[slide_index]

            # 既存の Table/Schedule（先に作った表など）を削除
            for shp in list(slide.shapes):
                name = getattr(shp, "name", "")
                # ここは Table*, Schedule*（大文字）だけ消すので、schedule1〜3 は消さない
                if name.startswith("Table") or name.startswith("Schedule"):
                    try:
                        slide.shapes._spTree.remove(shp._element)
                    except Exception:
                        pass

            # === schedule1 / schedule2 / schedule3 のプレースホルダ図形を取得 ===
            placeholders = {}
            for shp in slide.shapes:
                name = getattr(shp, "name", "")
                if name in ["schedule1", "schedule2", "schedule3"]:
                    placeholders[name] = shp

            # 3つともなくても動くようにする（ある分だけ使う）
            # DataFrameインデックスを整理
            df = calendar_df.reset_index(drop=True)
            total_rows = len(df)
            if total_rows == 0:
                return prs

            # 3ブロックに分割
            rows_per_block = math.ceil(total_rows / 3)

            # ===== カラー設定 =====
            header_fill_color   = RGBColor(230, 230, 230)  # ヘッダー：薄いグレー
            body_fill_color     = RGBColor(255, 255, 255)  # 平日：白
            holiday_fill_color  = RGBColor(240, 240, 240)  # 非営業日：さらに薄いグレー
            text_color          = RGBColor(0, 0, 0)        # 黒
            headers = ["日付", "曜日", "マイルストン"]

            # 各ブロック（1〜3）を、それぞれ schedule1〜3 の位置に描画
            for block_idx in range(3):
                start_idx = block_idx * rows_per_block
                end_idx = min(start_idx + rows_per_block, total_rows)
                block_df = df.iloc[start_idx:end_idx]

                if block_df.empty:
                    continue

                placeholder_name = f"schedule{block_idx + 1}"
                ph = placeholders.get(placeholder_name)
                if ph is None:
                    # schedule1/2/3 のどれかが無い場合、そのブロックはスキップ
                    continue

                # プレースホルダ図形の位置とサイズを取得
                left   = ph.left
                top    = ph.top
                width  = ph.width
                height = ph.height

                # プレースホルダを削除（同じ位置に表を置く）
                try:
                    slide.shapes._spTree.remove(ph._element)
                except Exception:
                    pass

                rows = len(block_df) + 1  # ヘッダー行 + データ行
                cols = 3

                table = slide.shapes.add_table(rows, cols, left, top, width, height).table

                # --- ヘッダー行 ---
                for j, h in enumerate(headers):
                    cell = table.cell(0, j)
                    cell.text = h

                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_fill_color

                    for p in cell.text_frame.paragraphs:
                        p.font.bold = True
                        p.font.size = Pt(12)
                        p.alignment = PP_ALIGN.CENTER
                        p.font.name = "Meiryo UI"
                        p.font.color.rgb = text_color

                # --- データ行 ---
                for i, (_, row) in enumerate(block_df.iterrows()):
                    table.cell(i + 1, 0).text = str(row.get("日付", ""))
                    table.cell(i + 1, 1).text = str(row.get("曜日", ""))
                    table.cell(i + 1, 2).text = str(row.get("マイルストン", ""))

                    is_holiday = bool(row.get("非営業日", False))

                    for j in range(3):
                        cell = table.cell(i + 1, j)

                        # 非営業日は薄いグレー、それ以外は白
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = holiday_fill_color if is_holiday else body_fill_color

                        for p in cell.text_frame.paragraphs:
                            p.font.size = Pt(11)
                            p.font.name = "Meiryo UI"
                            p.font.color.rgb = text_color

            return prs



        # ------------------------------------------------
        # Streamlit UI（ここで schedule_phase_draft を反映）
        # ------------------------------------------------
        st.subheader("主要な日程の入力")

        # ▼ オリエン内容の整理から抽出したマイルストンがあれば、それを優先して使用
        if "schedule_phase_draft" in st.session_state and st.session_state["schedule_phase_draft"]:
            base_phases = st.session_state["schedule_phase_draft"]
            df_phase = pd.DataFrame(base_phases)

            # duration_days がなければ追加（デフォルト3営業日）
            if "duration_days" not in df_phase.columns:
                df_phase["duration_days"] = 3
            df_phase["duration_days"] = df_phase["duration_days"].fillna(3).astype(int)

            # fixed_date は DateColumn で扱えるように date 型にしておく
            df_phase["fixed_date"] = pd.to_datetime(df_phase["fixed_date"], errors="coerce").dt.date

            st.info("『オリエン内容の整理』から抽出したマイルストンを初期値として読み込みました。必要に応じて工程名・営業日数を調整してください。")
        else:
            # フォールバック：従来のデフォルト工程
            default_phases = [
                {"name": "企画ご提案", "fixed_date": None, "duration_days": 2},
                {"name": "調査票案ご提示", "fixed_date": None, "duration_days": 3},
                {"name": "調査実施", "fixed_date": None, "duration_days": 5},
                {"name": "集計データアップ", "fixed_date": None, "duration_days": 2},
                {"name": "報告書提出", "fixed_date": None, "duration_days": 3},
            ]
            df_phase = pd.DataFrame(default_phases)

        edited_phases = st.data_editor(
            df_phase,
            hide_index=True,
            num_rows="dynamic",
            column_config={
                "name": st.column_config.TextColumn("工程名（マイルストン）"),
                "fixed_date": st.column_config.DateColumn("固定日（任意）"),
                "duration_days": st.column_config.NumberColumn("営業日数", min_value=1),
            },
            use_container_width=True,
            key="phases_editor_v7"
        )

        if "calendar_df_v7" not in st.session_state:
            st.session_state["calendar_df_v7"] = None

        if st.button("📅 スケジュールを生成", use_container_width=True):
            sched_df = compute_schedule_with_fixed(edited_phases)
            st.session_state["calendar_df_v7"] = make_calendar_table(sched_df)
            st.success("スケジュールを生成しました！")

        cal_df = st.session_state.get("calendar_df_v7")
        if cal_df is None:
            st.info("スケジュールが未生成です。上の『📅 スケジュールを生成』ボタンを押してください。")
        else:
            st.subheader("カレンダー表（編集可）")
            # 行数に応じて高さを自動調整
            row_height = 28  # 1行あたりの高さ(px)
            n_rows = len(cal_df)
            table_height = min(120 + row_height * n_rows, 900)  # 最大900pxまで拡張

            edited_cal = st.data_editor(
                cal_df.drop(columns=["非営業日"]),
                num_rows="fixed",
                use_container_width=True,
                height=table_height,
                column_config={
                    "日付": st.column_config.TextColumn("日付"),
                    "曜日": st.column_config.TextColumn("曜日"),
                    "マイルストン": st.column_config.TextColumn("マイルストン（編集可）"),
                },
                key="calendar_editor_v7"
            )

            st.markdown("---")

            # ★ スケジュール表 → スライド7 に反映
            if st.button("📤 このスケジュール表をスライド7に反映", use_container_width=True):
                pptx_path = st.session_state.get("pptx_path")
                if not pptx_path:
                    st.warning("PPTテンプレートを先にアップロードしてください。")
                else:
                    try:
                        prs = Presentation(pptx_path)
                        prs = reflect_schedule_to_slide7(prs, edited_cal)
                        prs.save(pptx_path)

                        st.success("スライド7（スケジュール案）にスケジュール表を反映しました！")

                        # 🔁 ここで再実行 → 冒頭のプレビューが更新後のPPTを読み込む
                        st.rerun()

                    except Exception as e:
                        st.error(f"PowerPoint反映中にエラーが発生しました: {e}")



        
    # =========================
    # 中央ペイン
    # === 概算見積（演算＆5パターン表示）===
    elif mode == "概算見積":
        st.markdown("## 概算見積")

        pptx_path = st.session_state.get("pptx_path")

        # ★ 前回反映フラグが立っていれば成功メッセージを一度だけ表示
        if st.session_state.get("estimate_applied"):
            st.success("スライド8（概算見積）の EDIT_amount1〜5 に反映しました！")
            st.session_state["estimate_applied"] = False

        # ---- PPTプレビュー表示（スライド8：画像プレビュー）----
        if pptx_path:
            try:
                images = pptx_to_images(pptx_path)
                if len(images) > 16:  # スライド8は index=7（0始まり）…テンプレ側に合わせて調整
                    st.image(images[16], caption="スライド8：概算見積", use_container_width=True)
                else:
                    st.warning("スライド8（概算見積）がテンプレートに存在しません。")
            except Exception as e:
                st.error(f"PPTプレビュー生成中にエラーが発生しました: {e}")
        else:
            st.info("PPTテンプレートをアップロードしてください。")

        st.markdown("---")
        st.markdown("### 🧮 入力内容にもとづく概算見積（右ペインで仕様を入力してください）")

        # ======================
        # 価格テーブル・関数群（既存ロジックをそのまま使用）
        # ======================
        import pandas as pd
        from pathlib import Path
        from datetime import datetime
        from pptx import Presentation
        from pptx.dml.color import RGBColor

        HOUR_RATE = 15000  # 人件費 1時間＝15,000円

        # 本調査価格表（代表値）
        MAIN_TABLE = {
            (20, 100): 187000,
            (40, 100): 328000,
            (20, 300): 216000,
            (40, 300): 382000,
            (20, 500): 255000,
            (40, 500): 460000,
            (20, 1000): 372000,
            (40, 1000): 675000,
        }

        # スクリーニング価格表（代表値）
        SCR_TABLE = {
            (5, 10000): 100000,
            (10, 10000): 180000,
            (15, 10000): 308000,
            (5, 20000): 130000,
            (10, 20000): 220000,
            (15, 20000): 358000,
            (5, 30000): 160000,
            (10, 30000): 260000,
            (15, 30000): 408000,
            (5, 50000): 220000,
            (10, 50000): 340000,
            (15, 50000): 508000,
            (5, 70000): 280000,
            (10, 70000): 420000,
            (15, 70000): 608000,
            (5, 100000): 370000,
            (10, 100000): 540000,
            (15, 100000): 758000,
        }

        def lookup_price(table: dict, q: int, n: int) -> int:
            """テーブルから一番近い組み合わせの価格をざっくり取得"""
            if not table:
                return 0

            best_key = None
            best_score = None
            for (tq, tn), price in table.items():
                dq = abs(tq - q)
                dn = abs(tn - n) / 1000  # サンプル差はスケール調整
                score = dq * dq + dn * dn
                if best_score is None or score < best_score:
                    best_score = score
                    best_key = (tq, tn)

            return table.get(best_key, 0)

        def to_man_yen(v: float) -> float:
            """円 → 万円"""
            return v / 10000.0

        # ======================
        # 右ペインで入力された値を session_state から取得
        # ======================
        hours_plan = float(st.session_state.get("hours_plan", 0.0))
        hours_field = float(st.session_state.get("hours_field", 0.0))
        hours_agg = float(st.session_state.get("hours_agg", 0.0))
        hours_analysis = float(st.session_state.get("hours_analysis", 0.0))

        scr_q = int(st.session_state.get("scr_q", 5))
        scr_n = int(st.session_state.get("scr_n", 10000))
        main_q = int(st.session_state.get("main_q", 20))
        main_n = int(st.session_state.get("main_n", 300))

        # 右ペイン未入力時のガード
        if hours_plan == hours_field == hours_agg == hours_analysis == 0 and \
        scr_q == 0 and scr_n == 0 and main_q == 0 and main_n == 0:
            st.info("右ペインで『企画費用（人件費）』と『実査費用（ベース仕様）』を入力してください。")
            st.stop()


        # ======================
        # 企画費用（人件費） 共通計算
        # ======================
        cost_plan = hours_plan * HOUR_RATE
        cost_field = hours_field * HOUR_RATE
        cost_agg = hours_agg * HOUR_RATE
        cost_analysis = hours_analysis * HOUR_RATE
        planning_total = cost_plan + cost_field + cost_agg + cost_analysis

        # スクリーニング費用（全パターン共通）
        scr_cost_base = lookup_price(SCR_TABLE, scr_q, scr_n)

        # ======================
        # 5パターンの仕様生成
        # ======================
        patterns = []

        def make_pattern(name: str, label: str, q: int, n: int):
            main_cost = lookup_price(MAIN_TABLE, q, n)
            survey_total = scr_cost_base + main_cost
            total_cost = planning_total + survey_total

            summary_lines = [
                f"■{label}",
                "",
                "【企画費用（人件費）】",
                f"・調査企画：{hours_plan:.1f}人時 ＝ {to_man_yen(cost_plan):,.1f} 万円",
                f"・調査実査：{hours_field:.1f}人時 ＝ {to_man_yen(cost_field):,.1f} 万円",
                f"・集計：{hours_agg:.1f}人時 ＝ {to_man_yen(cost_agg):,.1f} 万円",
                f"・分析・報告：{hours_analysis:.1f}人時 ＝ {to_man_yen(cost_analysis):,.1f} 万円",
                f"▶ 企画費用 小計：{to_man_yen(planning_total):,.1f} 万円",
                "",
                "【実査費用】",
                f"・スクリーニング：{scr_q}問 × {scr_n:,}ss ＝ {to_man_yen(scr_cost_base):,.1f} 万円",
                f"・本調査：{q}問 × {n:,}ss ＝ {to_man_yen(main_cost):,.1f} 万円",
                f"▶ 実査費用 小計：{to_man_yen(survey_total):,.1f} 万円",
                "",
                f"■概算合計：{to_man_yen(total_cost):,.1f} 万円（税別）",
            ]
            return {
                "name": name,
                "label": label,
                "q": q,
                "n": n,
                "main_cost": main_cost,
                "survey_total": survey_total,
                "total_cost": total_cost,
                "summary": "\n".join(summary_lines),
            }

        # 1) ベース仕様
        patterns.append(
            make_pattern(
                "pattern1",
                "ベース仕様（入力どおり）",
                main_q,
                main_n,
            )
        )

        # 2) 本調査サンプルサイズのみ半分
        patterns.append(
            make_pattern(
                "pattern2",
                "本調査サンプルサイズを半分にした場合",
                main_q,
                max(1, main_n // 2),
            )
        )

        # 3) 本調査サンプルサイズのみ2倍
        patterns.append(
            make_pattern(
                "pattern3",
                "本調査サンプルサイズを2倍にした場合",
                main_q,
                max(1, main_n * 2),
            )
        )

        # 4) 本調査質問数のみ5問減
        patterns.append(
            make_pattern(
                "pattern4",
                "本調査質問数を5問減らした場合",
                max(1, main_q - 5),
                main_n,
            )
        )

        # 5) 本調査質問数のみ5問増
        patterns.append(
            make_pattern(
                "pattern5",
                "本調査質問数を5問増やした場合",
                max(1, main_q + 5),
                main_n,
            )
        )

        # ======================
        # 5パターン概要テーブル表示
        # ======================
        st.markdown("### 📊 5パターンの比較サマリー")

        df_view = pd.DataFrame(
            [
                {
                    "パターン": p["label"],
                    "本調査質問数": p["q"],
                    "本調査サンプルサイズ": p["n"],
                    "概算合計（万円）": f"{to_man_yen(p['total_cost']):,.1f}",
                }
                for p in patterns
            ]
        )

        st.dataframe(df_view, use_container_width=True)

        st.markdown("---")
        st.markdown("### 📝 スライド貼り付けプレビュー（各パターン）")

        # テキストプレビュー＋ session_state に保存
        for idx, p in enumerate(patterns, start=1):
            key_txt = f"estimate_summary{idx}"
            st.session_state[key_txt] = p["summary"]

            st.markdown(f"#### パターン{idx}：{p['label']}")
            st.text_area(
                f"スライド用テキスト（EDIT_amount{idx} に反映）",
                value=p["summary"],
                height=260,
                key=f"estimate_summary_area_{idx}",
            )
            st.markdown("---")

        # ===============================
        # 📤 PowerPoint反映ボタン（EDIT_amount1〜5）
        # ===============================
        st.markdown("### 📤 5パターンを PowerPoint に反映（EDIT_amount1〜5）")

        if st.button("📤 5パターンすべてをスライド8に反映", use_container_width=True):
            pptx_path = st.session_state.get("pptx_path")

            if not pptx_path:
                st.warning("PPTテンプレートを先にアップロードしてください。")
            else:
                try:
                    from pptx import Presentation
                    from pathlib import Path
                    from datetime import datetime

                    prs = Presentation(pptx_path)
                    slide_index = 16  # スライド8（0始まり）

                    if slide_index < len(prs.slides):
                        slide = prs.slides[slide_index]

                        applied_count = 0
                        for idx in range(1, 6):
                            shape_name = f"EDIT_amount{idx}"
                            text_to_apply = st.session_state.get(f"estimate_summary{idx}", "")

                            if not text_to_apply:
                                continue

                            ok = set_text_to_named_shape(slide, shape_name, text_to_apply)

                            if ok:
                                shp = next((s for s in slide.shapes if s.name == shape_name), None)
                                if shp and getattr(shp, "has_text_frame", False):
                                    # ★ 概算見積だけフォントサイズ10ptに統一
                                    apply_text_format(shp, font_size=10)

                                # プレビュー用キャッシュ
                                st.session_state.edited_texts[shape_name] = text_to_apply
                                applied_count += 1
                            else:
                                st.warning(f"スライド8内に『{shape_name}』という名前のテキスト図形が見つかりませんでした。")

                        if applied_count > 0:
                            SLIDES_DIR = get_session_dir()
                            SLIDES_DIR.mkdir(parents=True, exist_ok=True)
                            out_path = SLIDES_DIR / f"estimate_slide8_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

                            prs.save(out_path)
                            st.session_state["pptx_path"] = str(out_path)

                            # プレビュー更新のためのフラグ
                            st.session_state["estimate_applied"] = True

                            st.success(f"スライド8（概算見積）に {applied_count} パターン分を反映しました！（フォントサイズ10pt）")
                            st.rerun()
                        else:
                            st.error("いずれの EDIT_amount1〜5 にもテキストを反映できませんでした。Shape名やテンプレート構成を確認してください。")

                    else:
                        st.error("スライド8がテンプレートに存在しません。")

                except Exception as e:
                    st.error(f"PowerPoint反映中にエラーが発生しました: {e}")



    # =========================
    # 中央ペイン
    # === パワーポイント出力 ===
    elif mode == "パワーポイントを出力":
        st.markdown("## 📤 PowerPoint出力")

        from pathlib import Path
        from datetime import datetime
        from pptx import Presentation

        pptx_path = st.session_state.get("pptx_path")

        if not pptx_path or not Path(pptx_path).is_file():
            st.warning(
                "まだPowerPointテンプレートがアップロードされていません。"
                "左ペインからテンプレートをアップロードしてください。"
            )
        else:
            st.caption(
                "現在のPPTビューアーに反映されている内容を、そのまま最終版PowerPointとして書き出します。"
                "（各モードでの『…をスライドXに反映』ボタンで更新された状態が保存されます）"
            )

            if st.button("💾 現在の内容で最終版PowerPointを作成", use_container_width=True):
                try:
                    # 現在の pptx_path の内容をそのまま別名保存
                    prs = Presentation(pptx_path)

                    SLIDES_DIR = get_session_dir()
                    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
                    out_path = SLIDES_DIR / f"proposal_final_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

                    prs.save(out_path)
                    st.session_state["final_pptx_path"] = str(out_path)

                    st.success(
                        "現在のPPTビューアーに反映されている内容をもとに "
                        "最終版PowerPointを作成しました。右ペインからダウンロードできます。"
                    )

                except Exception as e:
                    st.error(f"最終版PowerPoint作成中にエラーが発生しました: {e}")




# =========================
# =========================
# 右ペイン
# =========================
# =========================
with right:
    mode = st.session_state.get("selected_mode")

    # =========================
    # 右ペイン
    # === オリエン内容の整理 ===
    if mode == "オリエン内容の整理":
        st.subheader("オリエン内容の整理")
        st.caption("オリエン資料をもとに必要項目の下書きを作成します。")

        if st.button("下書き開始", use_container_width=True):
            ori_texts = "\n".join(st.session_state.get("uploaded_docs", []))

            if not ori_texts.strip():
                st.warning("オリエン資料をアップロードしてください。")
            else:
                with st.spinner("オリエン資料から項目を抽出中..."):
                    prompt = f"""
あなたは市場調査の専門家です。
以下のオリエン資料から以下のことをまとめてください。
特に言及がなければ項目ごとに「なし」と記載してください。

【出力形式】
・企業名：
・ブランド名：
・カテゴリー（市場）名：
・議事録の要約（500文字程度）：
・分析手法に関する要望：
・調査仕様に関する要望
    調査エリア：
    スクリーニング調査有無：
    対象者条件：
    質問数：
    サンプルサイズ：
    調査画面で画像や動画の提示：
    ウェイトバック集計の有無：
    自由回答のコーディング処理の有無：
    調査票作成（クライアントがやるか当社がやるか）：
    報告書は必要か：
・スケジュールに関する要望
    企画提案予定日：
    調査票や画像に関する提供可能日：
    希望する納期：
    請求日/月：
    クライアントの重要な会議日：
    その他スケジュールに関する要望：
・費用に関する要望
    見積金額上限：
    複数パターンの見積を希望しているか：
・会議参加者のお名前・役職・役割
・調査とは直接関係ないが雑談したこと：
・その他調査に関する特記事項（広告がいつから投下されるかなど）：


オリエン資料：
{ori_texts[:4000]}
"""
                    try:
                        response = client.chat.completions.create(
                            model=DEPLOYMENT,
                            messages=[
                                {"role": "system", "content": "あなたは市場調査の専門家です。"},
                                {"role": "user", "content": prompt},
                            ],
                            temperature=0.3,
                            max_tokens=900,  # ★長めに確保
                        )
                        ai_result = response.choices[0].message.content.strip()

                        # ★全文をセッションに保存（中央ペインで表示する用）
                        st.session_state["orien_outline_text"] = ai_result
                        st.session_state["orien_outline_editor"] = ai_result
                        # ★企業名だけも別途保存（表紙などと連携したい場合用）
                        # import re
                        # m = re.search(r"企業名[:：]\s*(.*)", ai_result)
                        # company = m.group(1).strip() if m else ""
                        # st.session_state["orien_company_text"] = company

                        st.success("オリエン内容の下書きを作成しました。中央ペインに表示します。")
                        st.rerun()

                    except Exception as e:
                        st.error(f"AI呼び出しエラー: {e}")


    # =========================
    # 右ペイン
    # ==== ブランド診断モード ====
    elif st.session_state.get("selected_mode") == "brand_diagnosis":
        st.subheader("カテゴリー・ブランド診断")
        st.caption("オリエン資料をもとにカテゴリー・ブランドを推測し、市場特性を検索します。")

        # --- 初期化 ---
        st.session_state.setdefault("target_category", "")
        st.session_state.setdefault("target_brand", "")

        ori_texts = "\n".join(st.session_state.get("uploaded_docs", []))

        # カテゴリー・ブランドを推測
        if st.button("📘 カテゴリー・ブランドを推測", use_container_width=True):
            if not ori_texts.strip():
                st.warning("オリエン資料をアップロードしてください。")
            else:
                with st.spinner("カテゴリーとブランドを推測中..."):
                    prompt = f"""
    あなたは市場調査の専門家です。
    以下のオリエン資料から、今回の調査対象となるカテゴリー（市場）とブランド名を推定してください。

    【出力形式】
    カテゴリー（市場）:
    ブランド:

    資料:
    {ori_texts[:4000]}
    """
                    try:
                        response = client.chat.completions.create(
                            model=DEPLOYMENT,
                            messages=[
                                {"role": "system", "content": "あなたは市場調査の専門家です。"},
                                {"role": "user", "content": prompt},
                            ],
                            temperature=0.5,
                            max_tokens=200,
                        )
                        ai_result = response.choices[0].message.content

                        import re
                        cat_match = re.search(r"カテゴリー（市場）[:：]\s*(.*)", ai_result)
                        brand_match = re.search(r"ブランド[:：]\s*(.*)", ai_result)

                        st.session_state["target_category"] = cat_match.group(1).strip() if cat_match else ""
                        st.session_state["target_brand"] = brand_match.group(1).strip() if brand_match else ""

                        st.success("カテゴリーとブランドを抽出しました。下の欄で確認・編集できます。")
                        st.rerun()

                    except Exception as e:
                        st.error(f"AI呼び出しエラー: {e}")


        # 手動編集欄
        st.text_input(
            "対象カテゴリー（市場）",
            key="target_category",
            placeholder="例：清涼飲料、化粧品、通信キャリアなど",
        )
        st.text_input(
            "対象ブランド",
            key="target_brand",
            placeholder="例：キッザニア、SUUMO、カローラ など",
        )

        st.divider()


        # カテゴリー・ブランドについて検索
        st.markdown("カテゴリー・ブランドについて検索")

        if st.button("カテゴリー・ブランドについて検索", use_container_width=True):
            cat = st.session_state.get("target_category", "")
            brand = st.session_state.get("target_brand", "")
            if not cat:
                st.warning("カテゴリーを入力してください。")
            else:
                with st.spinner("市場特性を検索中..."):
                    prompt = f"""
    あなたは市場分析の専門家です。
    次のカテゴリーとブランドに関する市場構造と消費行動特性を整理してください。

    【カテゴリー】{cat}
    【ブランド】{brand}

    出力は以下の2表形式のMarkdownで記載してください。
    # カテゴリーに関する検索項目
    |項目|内容|
    |市場タイプ|FMCG／耐久財／サービス材／BtoB／公共／非営利／デジタルプロダクト|
    |市場成長ステージ|成長／成熟／停滞／衰退／新興|
    |市場競争構造|リーダー1強／寡占2〜3者／分散|
    |ブランド特性・ポジション|リーダー／チャレンジャー／フォロワー／ニッチ／新規参入|
    |購買・意思決定構造|高関与／低関与／集団意思決定／専門家介在／衝動購買|
    |顧客心理構造|感情重視型／機能重視型／信頼重視型など|
    |流通・販売構造|店頭中心／EC中心／直販／代理店など|
    |顧客関係構造|一回購入型／サブスク／リピート中心／契約継続型|
    |組織・ブランド構造|単一ブランド／マルチブランド|
    |社会・文化的文脈|ライフスタイルトレンド／社会課題との接点など|

    # カテゴリーの消費行動特性
    |項目|内容|
    |検討期間|長期／短期／反復購入|
    |情報収集経路|SNS／Web／来店／紹介など|
    |購入決定単位|個人／家族／グループ|
    |再購入／継続構造|定期購入／都度購入|

    出力は以下のような**区切り線を含まないシンプルな表形式（縦線とセルだけ）**で記載してください。
    区切り線（---）やMarkdownヘッダー構文は入れないでください。



"""
                    try:
                        response = client.chat.completions.create(
                            model=DEPLOYMENT,
                            messages=[
                                {"role": "system", "content": "あなたは市場分析の専門家です。"},
                                {"role": "user", "content": prompt},
                            ],
                            temperature=0.6,
                            max_tokens=900,
                        )
                        result = response.choices[0].message.content

                        import pandas as pd, re

                        def extract_md_table(md_text, header):
                            if header in md_text:
                                section = md_text.split(header, 1)[1]
                                table_part = section.split("#")[0]
                                rows = [
                                    ln.strip()
                                    for ln in table_part.splitlines()
                                    if "|" in ln and not ln.startswith("|項目|----|")
                                ]
                                data = []
                                for ln in rows:
                                    cols = [c.strip() for c in ln.strip("|").split("|")]
                                    if len(cols) >= 2:
                                        data.append(cols[:2])
                                if data:
                                    df = pd.DataFrame(data[1:], columns=data[0])
                                    return df
                            return pd.DataFrame(columns=["項目", "内容"])

                        st.session_state["df_category_structure"] = extract_md_table(result, "# カテゴリーに関する検索項目")
                        st.session_state["df_behavior_traits"] = extract_md_table(result, "# カテゴリーの消費行動特性")

                        st.success("市場特性を整理しました。中央ペインに表示されます。")

                        # ---- 追加：マーケティングファネル生成 ----
                        with st.spinner("マーケティングファネルを生成中..."):
                            prompt_funnel = f"""
あなたはブランドマーケティングの専門家であり、人の思考を支援するアシスタントです。
以下のカテゴリーとブランドについて、消費者が「認知」から「再接点・ロイヤリティ」に至るまでの
マーケティングファネルをツリー構造で整理してください。

【出力形式】
- 1階層目：ファネル段階（認知→興味→検討→購入→再接点・ロイヤリティ）
- 2階層目：それぞれの段階における「トリガー（促進要因）」と「障壁（阻害要因）」
- 3階層目：各トリガー・障壁の具体例（2〜3項目ずつ、短文または名詞句）
- 出力は階層インデント（記号やスペース）で明確にしてください。
- Markdownの箇条書き（- や *）を使って構造を表現してください。
- 余計な説明文や序文は不要です。ツリー構造のみを出力してください。

【カテゴリー】{cat}
【ブランド】{brand}

出力例：

- 認知（Awareness）
  - トリガー
    - 広告露出（SNS・Web・業界誌）
    - 口コミ・レビュー
  - 障壁
    - 認知不足
    - 競合の露出優位
- 興味・関心（Interest）
  - トリガー
    - 成功事例の紹介
    - 無料体験の提供
  - 障壁
    - 情報過多による混乱
    - 利用メリットが伝わらない
...
    """
                            response_funnel = client.chat.completions.create(
                                model=DEPLOYMENT,
                                messages=[
                                    {"role": "system", "content": "あなたはブランドマーケティングの専門家です。"},
                                    {"role": "user", "content": prompt_funnel},
                                ],
                                temperature=0.6,
                                max_tokens=1800,
                            )
                            st.session_state["funnel_text"] = response_funnel.choices[0].message.content

                        st.success("マーケティングファネルを生成しました。中央ペインに表示されます。")
                        st.rerun()

                    except Exception as e:
                        st.error(f"AI呼び出しエラー: {e}")

    # =========================
    # 右ペイン
    # === 表紙 ===
    # elif mode == "表紙":
    #     st.subheader("右ペイン：表紙操作")
    #     if st.button("💡 ダミーボタン（表紙）", use_container_width=True):
    #         st.session_state["message_center"] = "💬 『表紙』でダミーボタンが押されました。"
    #         st.rerun()


    # =========================
    # 右ペイン
    # === キックオフノート ===
    elif mode == "キックオフノート":
        st.subheader("キックオフノート")
        st.caption("オリエン資料とブランド診断結果をもとに、①〜⑥の下書き生成します。")

        # ------------------------------------------------------------
        # 調査目的のマトリクス選択
        # ------------------------------------------------------------
        PURPOSE_MATRIX = {
"市場・競合把握": "市場規模、成長性、競合構造などの理解を目的とした調査です。", 
"ニーズボリューム把握": "消費者ニーズの量的分布を明らかにし、優先ターゲットを特定します。", 
"実態・意識把握": "消費者の行動実態や意識構造を明らかにする定性・定量調査です。", 
"ニーズ探索": "潜在的な消費者ニーズやウォンツを発掘・探索します。", 
"アイデアスクリーニング": "複数のアイデア案を評価・選抜するための初期テストを行います。", 
"コンセプト受容性把握": "商品・サービスコンセプトの受容度、共感度、理解度を測定します。", 
"スぺック評価把握": "製品スペック（機能・特徴）の重要度や評価ポイントを明らかにします。", 
"価格弾力性把握": "価格設定に対する需要反応（価格弾力性）を推定します。", 
"需要予測": "市場シェアや販売量の見込みを予測する調査です。", 
"訴求ポイント把握": "広告・コミュニケーションで強調すべきメッセージを明確化します。", 
"浸透状況把握": "ブランド・製品の市場浸透率や認知度、利用率を測定します。", 
"サービス使用評価": "実際のサービス利用体験を通じた満足度・課題を抽出します。", 
"プロモ効果測定": "キャンペーンやプロモーションの効果を定量的に評価します。", 
"ユーザー評価": "既存ユーザーからの製品・サービス評価を把握します。"
        }

        selected_purpose = st.selectbox(
            "◆調査テーマを選択してください。",
            list(PURPOSE_MATRIX.keys()),
            key="kickoff_selected_purpose"
        )

        st.divider()

        # ------------------------------------------------------------
        # 🪄 AI下書き生成（①〜⑥）
        # ------------------------------------------------------------
        if st.button("下書きを生成", use_container_width=True):
            ori_texts = "\n".join(st.session_state.get("uploaded_docs", []))
            orien_outline_text = st.session_state.get("orien_outline_text", "")
            cat_df = st.session_state.get("df_category_structure")
            beh_df = st.session_state.get("df_behavior_traits")
            funnel_text = st.session_state.get("funnel_text", "")

            if not ori_texts.strip():
                st.warning("オリエン資料をアップロードしてください。")
            else:
                with st.spinner("キックオフノートの下書きを作成中..."):
                    matrix_text = PURPOSE_MATRIX.get(selected_purpose, "")
                    cat_text = cat_df.to_markdown(index=False) if cat_df is not None and not cat_df.empty else ""
                    beh_text = beh_df.to_markdown(index=False) if beh_df is not None and not beh_df.empty else ""

                    prompt = f"""
    あなたは市場調査設計の専門家です。
    以下のオリエン資料、ブランド診断結果、調査目的マトリクスをもとに、
    調査設計の初期段階で用いる「キックオフノート」を作成してください。

    【出力形式】
    【目標】
    【現状】
    【ビジネス課題】
    【調査目的】
    【問い】
    【仮説】
    【ポイント】

    
    【条件】
    - 各項目は80〜120字以内
    - オリエン資料にある固有名詞や文脈を十分に生かしてください。
    - 【目標】や【現状】は経営課題や社会問題など、調査では解決できない抽象課題は避けてください。
      あくまで「消費者・市場・ブランド・広告・顧客体験」など、市場調査で仮説検証できる範囲に課題を限定してください。
    - 【問い】はオリエンシートやブランド診断を踏まえた現在の対象ブランドの"リサーチクエスチョン"のことです。
      ブランド全体について問う場合と広告やプロダクト/サービス、顧客接点など施策について問う場合があります。
    - 【ポイント】にはなぜキックオフノートの各項目にそう記載したのか、特に注意すべき点や補足説明を簡潔に記載してください。
      
    【入力データ】
    ▼オリエン内容の整理（抜粋）
    {orien_outline_text[:2000]}

    ▼ブランド診断：カテゴリー構造
    {cat_text}

    ▼ブランド診断：消費行動特性
    {beh_text}

    ▼マーケティングファネル
    {funnel_text}

    ▼選択した調査目的
    {selected_purpose}：{matrix_text}


    【禁止事項】
    - ###、** などの記号は使わないでください。
    """

                    try:
                        response = client.chat.completions.create(
                            model=DEPLOYMENT,
                            messages=[
                                {"role": "system", "content": "あなたは市場調査設計の専門家です。"},
                                {"role": "user", "content": prompt},
                            ],
                            temperature=0.6,
                            max_tokens=900,
                        )

                        result = response.choices[0].message.content
                        sections = parse_ai_output(result)

                        # セッションに保存
                        for key in sections:
                            st.session_state[f"ai_{key}"] = sections[key]

                        st.success("キックオフノートの下書きを生成しました！中央ペインに反映されます。")
                        st.rerun()

                    except Exception as e:
                        st.error(f"AI呼び出しエラー: {e}")


    # =========================
    # 右ペイン
    # === 問いの分解 ===
    elif mode == "問いの分解":
        st.subheader("問いの分解")
        st.caption("『問い』を検証するためのサブクエスチョンを生成します。")

        if st.button("下書きを生成", use_container_width=True):
            ori_texts = "\n".join(st.session_state.get("uploaded_docs", []))
            orien_outline_text = st.session_state.get("orien_outline_text", "")
            cat_df = st.session_state.get("df_category_structure")
            beh_df = st.session_state.get("df_behavior_traits")
            main_question = st.session_state.get("ai_問い", "")

            if not main_question.strip():
                st.warning("キックオフノート⑤『問い』が生成または入力されていません。")
            elif not ori_texts.strip():
                st.warning("オリエン資料をアップロードしてください。")
            else:
                with st.spinner("サブクエスチョンとアンケート項目を検討中..."):
                    cat_text = (
                        cat_df.to_markdown(index=False)
                        if cat_df is not None and not cat_df.empty
                        else ""
                    )
                    beh_text = (
                        beh_df.to_markdown(index=False)
                        if beh_df is not None and not beh_df.empty
                        else ""
                    )

                    prompt = f"""
あなたは市場調査設計の専門家です。
以下の情報をもとに、キックオフノート⑤『問い』（リサーチクエスチョン）を深掘りするための
【サブクエスチョン】を提案してください。サブクエスチョンへのアプローチは分析軸、評価項目、主な分析アプローチ、
読み方・示唆例を含めて具体的に示してください。

クロス集計分析の場合の例を示します。
【出力形式】
- サブクエスチョン1：認知度に影響を与える要因は何か？   
  - 分析軸：性年代など
  - 評価項目：認知度、利用率など
  - 主な分析アプローチ：性年代ごとに認知度の違いを比較する
  - 読み方・示唆例：若年層で認知度が低い場合、若年層向けの広告強化が必要など
- サブクエスチョン2：購入者タイプごとに主に利用する情報源は何か？
  - 分析軸：ヘビー層、ライト層など
  - 評価項目：購入タイプ
  - 主な分析アプローチ：ヘビー層ライト層ごとに情報源の違いを比較する
  - 読み方・示唆例：ヘビー層はSNS、ライト層は店頭広告が主な情報源など


【キックオフノート⑤ 問い】
{main_question}

▼オリエン内容の整理（抜粋）
 {orien_outline_text[:2000]}

【ブランド診断：カテゴリー構造】
{cat_text}

【ブランド診断：消費行動特性】
{beh_text}

【禁止事項】
 - ###、** などの記号は使わないでください。
 - サブクエスチョンはどの問い（リサーチクエスチョン）にも対応しているのかが分かるように具体的に記載してください。
 - 1つの問いに対して、最大3つのサブクエスチョンを提案してください。
"""

                    try:
                        response = client.chat.completions.create(
                            model=DEPLOYMENT,
                            messages=[
                                {"role": "system", "content": "あなたは市場調査設計の専門家です。"},
                                {"role": "user", "content": prompt},
                            ],
                            temperature=0.6,
                            max_tokens=2000,
                        )
                        ai_text = response.choices[0].message.content

                        # ★ 生テキストを保存（中央ペインのテキストエリア用）
                        st.session_state["ai_subquestions"] = ai_text

                        # ★ パースして構造化データも保存（問いの分解ビュー & 分析アプローチ用）
                        st.session_state["subq_list"] = parse_subquestions(ai_text)

                        st.success("下書きを生成しました！中央ペインおよび分析アプローチで利用できます。")
                        st.rerun()

                    except Exception as e:
                        st.error(f"AI呼び出しエラー: {e}")


    # =========================
    # 右ペイン
    # === 分析アプローチ ===
    elif mode == "分析アプローチ":
        st.subheader("分析アプローチ")
        st.caption("『問いの分解』で作成したサブクエスチョンを、分析アプローチ視点で項目に整理します。")

        # 『問いの分解』で保存した構造化データ
        subq_list = st.session_state.get("subq_list", [])

        if not subq_list:
            st.info("先に『問いの分解』モードでサブクエスチョンを生成してください。")
        else:

            # 🔽 ここから新機能：AIで6項目に分解した下書きを作成
            if st.button("下書きを作成", use_container_width=True):
                ori_texts = "\n".join(st.session_state.get("uploaded_docs", []))
                orien_outline_text = st.session_state.get("orien_outline_text", "")
                cat_df = st.session_state.get("df_category_structure")
                beh_df = st.session_state.get("df_behavior_traits")

                kickoff = {
                    "目標": st.session_state.get("ai_目標", ""),
                    "現状": st.session_state.get("ai_現状", ""),
                    "ビジネス課題": st.session_state.get("ai_ビジネス課題", ""),
                    "調査目的": st.session_state.get("ai_調査目的", ""),
                    "問い": st.session_state.get("ai_問い", ""),
                    "仮説": st.session_state.get("ai_仮説", ""),
                }

                # サブクエスチョン一覧（AIに渡す用）
                subq_text_lines = []
                for i, sq in enumerate(subq_list, 1):
                    subq_text_lines.append(f"SQ{i}: {sq.get('subq', '')}")
                subq_text = "\n".join(subq_text_lines)

                # 参考情報
                cat_text = cat_df.to_markdown(index=False) if cat_df is not None and not cat_df.empty else ""
                beh_text = beh_df.to_markdown(index=False) if beh_df is not None and not beh_df.empty else ""

                import json

                with st.spinner("サブクエスチョンごとの分析アプローチ案を検討中..."):
                    prompt = f"""
あなたは市場調査設計の専門家です。
以下のサブクエスチョンそれぞれについて、次の6項目の観点から分析アプローチの下書きを作成してください。

【対象となる6項目】
- id: "SQ1" のようなID
- subq: サブクエスチョン本文
- axis: 分析軸（セグメント）
- metric: 評価項目
- approach: 主な分析アプローチ（どのような切り口で分析するか）
- hypothesis: 検証する仮説（どのような結果が出ると何が言えるのか）

▼オリエン内容の整理（抜粋）
 {orien_outline_text[:2000]}

▼ブランド診断：カテゴリー構造
{cat_text}

▼ブランド診断：消費行動特性
{beh_text}

▼キックオフノート
{kickoff}

【サブクエスチョン一覧】
{subq_text}

【出力形式】
- 必ず JSON 配列のみを出力してください（余計な文章やコードブロックは書かないこと）
- 形式の例：

[
  {{
    "id": "SQ1",
    "subq": "・・・",
    "axis": "・・・",
    "metric": "・・・",
    "approach": "・・・",
    "hypothesis": "・・・"
  }},
  {{
    "id": "SQ2",
    "subq": "・・・",
    "axis": "・・・",
    "metric": "・・・",
    "approach": "・・・",
    "hypothesis": "・・・"
  }}
]

- 配列の要素数は、入力されたサブクエスチョンの数と同じにしてください。
- axis: 分析軸（セグメント）の案が複数ある場合は最も優先度の高いもの1つを提示してください。　
  また、分析軸案の後に（）で具体的な項目を記載してください。
- metric: 評価項目の案が複数ある場合は最も重要なもの1つを提示してください。
- metric: 評価項目案の後に（）で具体的な項目を記載してください。
  例：評価指標の場合は（あてはまる、ややあてはまる）など尺度の項目、イメージ項目の場合は（自分らしい、新しい）など
- approach: 主な分析アプローチ（どのような切り口で分析するか）は、以下の形式で記載してください。
  例：「性年代ごとに認知度の違いを比較する」「購入タイプ別に情報源の違いを分析する」など
- hypothesis: 検証する仮説（どのような結果が出ると何が言えるのか）の語尾に「～の可能性が高い（ある）」を用いないでください。
"""

                    try:
                        response = client.chat.completions.create(
                            model=DEPLOYMENT,
                            messages=[
                                {"role": "system", "content": "あなたは市場調査設計の専門家です。"},
                                {"role": "user", "content": prompt},
                            ],
                            temperature=0.6,
                            max_tokens=2000,
                        )
                        ai_text = response.choices[0].message.content.strip()

                        # ```json ... ``` で返ってきた場合のガード
                        if ai_text.startswith("```"):
                            ai_text = ai_text.strip("`")
                            ai_text = ai_text.replace("json", "", 1).strip()

                        try:
                            blocks = json.loads(ai_text)
                            if not isinstance(blocks, list):
                                raise ValueError("JSON配列ではありません。")

                        except Exception:
                            st.error("AI出力をJSON配列として解釈できませんでした。出力内容を確認してください。")
                            st.code(ai_text)
                        else:
                            # セッションに保存：中央ペインで参照する
                            st.session_state["analysis_blocks"] = blocks
                            # 以前の表示テキストもリセットしておく
                            if "analysis_block_texts" in st.session_state:
                                del st.session_state["analysis_block_texts"]

                            st.success("サブクエスチョン別の分析アプローチ案を作成しました。中央ペインに表示します。")
                            st.rerun()

                    except Exception as e:
                        st.error(f"AI呼び出しエラー: {e}")



    # =========================
    # 右ペイン
    # === 対象者条件を検討 ===
    elif mode == "対象者条件を検討":
        st.subheader("対象者条件を検討")
        st.caption("オリエン資料・ブランド診断・キックオフノート・問い分解の内容をもとに対象者条件を提案します。")

        if st.button("下書きを作成", use_container_width=True):
            ori_texts = "\n".join(st.session_state.get("uploaded_docs", []))
            orien_outline_text = st.session_state.get("orien_outline_text", "")
            cat_df = st.session_state.get("df_category_structure")
            beh_df = st.session_state.get("df_behavior_traits")
            main_question = st.session_state.get("ai_問い", "")
            subquestions = st.session_state.get("ai_subquestions", "")
            kickoff = {
                "目標": st.session_state.get("ai_目標", ""),
                "現状": st.session_state.get("ai_現状", ""),
                "ビジネス課題": st.session_state.get("ai_ビジネス課題", ""),
                "調査目的": st.session_state.get("ai_調査目的", ""),
                "問い": st.session_state.get("ai_問い", ""),
                "仮説": st.session_state.get("ai_仮説", ""),
            }

            if not ori_texts.strip():
                st.warning("オリエン資料をアップロードしてください。")
            else:
                with st.spinner("調査対象者条件を検討中..."):
                    cat_text = cat_df.to_markdown(index=False) if cat_df is not None and not cat_df.empty else ""
                    beh_text = beh_df.to_markdown(index=False) if beh_df is not None and not beh_df.empty else ""

                    prompt = f"""
    あなたは市場調査設計の専門家です。
    以下の情報をもとに、この調査の「対象者条件」を検討してください。

    【出力形式】
    - 対象者イメージ：　※1行で簡潔に記載してください。
    - 地域条件：
    - 年齢・性別条件：
    - 属性・利用行動条件：
    - 除外条件：

    【オリエン内容の整理（抜粋）】
    {orien_outline_text[:2000]}

    【ブランド診断：カテゴリー構造】
    {cat_text}

    【ブランド診断：消費行動特性】
    {beh_text}

    【キックオフノート】
    {kickoff}

    【問いの分解（AI生成サブクエスチョン）】
    {subquestions}

    - 条件は、全国／20–69歳男女／該当カテゴリー利用者などの一般的なフォーマットを基本に、
      調査目的との整合性を意識して作成してください。
    - 対象者イメージは冒頭に簡潔に記載してください（例：20〜30代女性のヘビーユーザーなど）。
    - 表記に**などの記号は使わないでください。
    - 「# 対象者条件案」など冒頭の見出しも不要です。
    - 「補足」や「説明文」も不要です。
    """

                    try:
                        response = client.chat.completions.create(
                            model=DEPLOYMENT,
                            messages=[
                                {"role": "system", "content": "あなたは市場調査設計の専門家です。"},
                                {"role": "user", "content": prompt},
                            ],
                            temperature=0.6,
                            max_tokens=500,
                        )
                        ai_text = response.choices[0].message.content.strip()

                        st.session_state["ai_target_condition"] = ai_text
                        st.success("調査対象者条件を生成しました！中央ペインに反映されます。")
                        st.rerun()

                    except Exception as e:
                        st.error(f"AI呼び出しエラー: {e}")




    # =========================
    # 右ペイン
    # === 調査項目案 ===
    elif mode == "調査項目案":
        st.subheader("調査項目案")
        st.caption("調査項目案を作成します。")

        if st.button("下書きを作成", use_container_width=True):
            # 「オリエン内容の整理」で作成したテキストを参照
            orien_outline_text = st.session_state.get("orien_outline_text", "")

            cat_df = st.session_state.get("df_category_structure")
            beh_df = st.session_state.get("df_behavior_traits")
            kickoff = {
                "目標": st.session_state.get("ai_目標", ""),
                "現状": st.session_state.get("ai_現状", ""),
                "ビジネス課題": st.session_state.get("ai_ビジネス課題", ""),
                "調査目的": st.session_state.get("ai_調査目的", ""),
                "問い": st.session_state.get("ai_問い", ""),
                "仮説": st.session_state.get("ai_仮説", ""),
            }
            subquestions = st.session_state.get("ai_subquestions", "")
            target_condition = st.session_state.get("ai_target_condition", "")

            if not orien_outline_text.strip():
                st.warning("先に『オリエン内容の整理』で下書きを作成してください。")
            else:
                with st.spinner("調査項目案を検討中..."):
                    cat_text = cat_df.to_markdown(index=False) if cat_df is not None and not cat_df.empty else ""
                    beh_text = beh_df.to_markdown(index=False) if beh_df is not None and not beh_df.empty else ""

                    prompt = f"""
    あなたは市場調査設計の専門家です。
    以下の情報をもとに、この調査で実施すべき調査項目案を提案してください。

    【出力条件】
    - 選択肢は不要（設問文のみ）
    - 設問文は質問文形式でなく、調査項目名として簡潔に表現する
      例：過去3年以内にキッザニアを訪れた経験はありますか？の場合、「キッザニア訪問経験」など
    - 各バージョンで網羅性と実務的な順序を意識する
    - 各バージョンは下記の見出しごとに分けて出力する
    - 各バージョンで「ちょうど」下記の問数になるようにする（10問／20問／30問／40問）
    - 各設問は1行で簡潔に（目安：全角60文字以内）
    - 見出しと設問リスト以外の説明文は出力しない

    【出力形式】
    # 10問バージョン
    1. ...
    2. ...
    （10問まで）

    # 20問バージョン
    1. ...
    （20問まで）

    # 30問バージョン
    1. ...
    （30問まで）

    # 40問バージョン
    1. ...
    （40問まで）

    【オリエン内容の整理（抜粋）】
    {orien_outline_text[:2000]}

    【ブランド診断：カテゴリー構造】
    {cat_text}

    【ブランド診断：消費行動特性】
    {beh_text}

    【キックオフノート】
    {kickoff}

    【問いの要因分解】
    {subquestions}

    【対象者条件】
    {target_condition}

   """

                    try:
                        response = client.chat.completions.create(
                            model=DEPLOYMENT,
                            messages=[
                                {"role": "system", "content": "あなたは市場調査設計の専門家です。"},
                                {"role": "user", "content": prompt},
                            ],
                            temperature=0.6,
                            max_tokens=3200,  # かなり余裕を持たせる
                        )
                        ai_text = response.choices[0].message.content.strip()

                        # デバッグ用に生テキストも一応保存しておくと便利
                        st.session_state["ai_survey_items_raw"] = ai_text

                        # ---- 出力を分割して辞書に格納 ----
                        import re
                        versions = {}
                        for ver in ["10問", "20問", "30問", "40問"]:
                            pattern = rf"#\s*{ver}バージョン(.*?)(?=#\s*\d+問バージョン|$)"
                            m = re.search(pattern, ai_text, re.DOTALL)
                            versions[ver] = m.group(1).strip() if m else ""

                        st.session_state["ai_survey_items"] = versions
                        st.success("調査項目案を生成しました！中央ペインに反映されます。")
                        st.rerun()

                    except Exception as e:
                        st.error(f"AI呼び出しエラー: {e}")




    # =========================
    # 右ペイン
    # === 調査仕様案 ===
    elif mode == "調査仕様案":
        st.subheader("調査仕様案")
        st.caption("『調査仕様の下書きを作成します。")

        if st.button("下書きを作成", use_container_width=True):
            # 入力ソースを取得
            orien_outline_text = st.session_state.get("orien_outline_text", "")
            target_condition = st.session_state.get("ai_target_condition", "")
            survey_items_selected = st.session_state.edited_texts.get("EDIT1", "")

            if not orien_outline_text.strip():
                st.warning("先に『オリエン内容の整理』で下書きを作成してください。")
            else:
                with st.spinner("調査仕様の下書きを作成中..."):
                    cat_df = st.session_state.get("df_category_structure")
                    beh_df = st.session_state.get("df_behavior_traits")

                    cat_text = cat_df.to_markdown(index=False) if cat_df is not None and not cat_df.empty else ""
                    beh_text = beh_df.to_markdown(index=False) if beh_df is not None and not beh_df.empty else ""

                    # JSON形式で返すように指示してパースしやすくする
                    import json

                    prompt = f"""
    あなたは市場調査設計の専門家です。
    以下の情報をもとに、この調査の「調査仕様案」を項目ごとに整理してください。

    【入力情報】
    ▼オリエン内容の整理
    {orien_outline_text[:2000]}

    ▼対象者条件
    {target_condition}

    ▼調査項目案（採用版：PPT EDIT1に反映した内容）
    {survey_items_selected}

    ▼参考情報：カテゴリー構造
    {cat_text}

    ▼参考情報：消費行動特性
    {beh_text}

    【出力する項目】
    - 調査手法
    - 抽出方法
    - 調査地域
    - 対象者条件
    - サンプルサイズ
    - 調査ボリューム
    - 提示物
    - 集計・分析仕様
    - 自由回答データの処理
    - 業務範囲
    - 納品物
    - インスペクションの方法
    - 謝礼の種類
    - 備考

    【出力形式】
    次のキーを持つ JSON オブジェクト「だけ」を出力してください。
    余計な説明文やコードブロック（```）は出力しないでください。

    {{
      "調査手法": "...",
      "抽出方法": "...",
      "調査地域": "...",
      "対象者条件": "...",
      "サンプルサイズ": "...",
      "調査ボリューム": "...",
      "提示物": "...",
      "集計・分析仕様": "...",
      "自由回答データの処理": "...",
      "業務範囲": "...",
      "納品物": "...",
      "インスペクションの方法": "...",
      "謝礼の種類": "...",
      "備考": "..."
    }}

    - 調査手法は特に明記がなければ「インターネット調査」を基本としてください。
      対象者条件の検討の中で、属性以外の条件がある場合は（スクリーニングあり）と付記してください。
    - 抽出方法は特に明記がなければ「割付抽出」としてください。
    - 対象者条件は、前述の対象者条件案を参考に、調査仕様として適切な形式に整えてください。
    - 調査ボリュームはスクリーニング調査と本調査を2行に分けて記載してください。
      本調査のボリュームは、調査項目案の選択結果を記載してください。
    - 自由回答データの処理は、オリエン内容のテキストに記載がなければ「なし」を基本としてください。
    - インスペクションの方法は、オリエン内容のテキストに記載がなければ「性別・年齢（2歳以上）のアンマッチの場合は、対象除外とする。」を基本としてください。
    - 謝礼の種類は、オリエン内容のテキストに記載がなければ「ポイント謝礼」を基本としてください。
    """

                    try:
                        response = client.chat.completions.create(
                            model=DEPLOYMENT,
                            messages=[
                                {"role": "system", "content": "あなたは市場調査設計の専門家です。"},
                                {"role": "user", "content": prompt},
                            ],
                            temperature=0.5,
                            max_tokens=1000,
                        )

                        ai_text = response.choices[0].message.content.strip()

                        # 念のため ```json ... ``` で返ってきた場合も対応
                        if ai_text.startswith("```"):
                            ai_text = ai_text.strip("`")
                            ai_text = ai_text.replace("json", "", 1).strip()

                        try:
                            spec_obj = json.loads(ai_text)
                        except Exception:
                            st.error("AI出力をJSONとして解釈できませんでした。出力内容を確認してください。")
                            st.code(ai_text)
                        else:
                            # SPEC_ITEMS に従って session_state に保存
                            for label, key in SPEC_ITEMS:
                                st.session_state[key] = spec_obj.get(label, "")

                            st.success("調査仕様の下書きを作成しました。中央ペインに表示します。")
                            st.rerun()

                    except Exception as e:
                        st.error(f"AI呼び出しエラー: {e}")


    # =========================
    # 右ペイン
    # === スケジュール案 ===
    elif mode == "スケジュール案":
        st.subheader("スケジュール案")
        st.caption("下書きを作成します。")

        # ▼ オリエン内容の整理テキストを取得
        orien_outline_text = st.session_state.get("orien_outline_text", "")

        # ▼ 下書き作成ボタン
        if st.button("下書きを作成", use_container_width=True):
            if not orien_outline_text.strip():
                st.warning("先に『オリエン内容の整理』で下書きを作成してください。")
            else:
                with st.spinner("オリエン内容からマイルストン案を抽出中..."):
                    import json

                    # オリエン整理テキストの「スケジュールに関する要望」部分から
                    # マイルストン名と固定日（ある場合）をJSON配列で返すようにAIに指示
                    prompt = f"""
あなたは市場調査プロジェクトのプロジェクトマネージャーです。
以下の「オリエン内容の整理」テキストの中から、スケジュールに関する項目と日付情報を整理してください。

【入力テキスト（オリエン内容の整理）】
{orien_outline_text[:2000]}

特に、次のような項目を優先して確認してください：
- 企画提案予定日
- ご発注予定日
- 調査票案受領予定日 ※お客様が調査票を作成する場合
- インテージから調査票送付日
- 調査開始日
- 調査終了日
- データ納品日
- 報告書の納品日 ※業務範囲に報告書納品がある場合

【出力条件】
- 出力は JSON 配列「だけ」としてください（説明文やコードブロックは不要）
- 配列の各要素は以下のキーを持つオブジェクトとします

[
  {{
    "name": "企画提案予定日",
    "fixed_date": "2025-02-10"
  }},
  {{
    "name": "調査票や画像の提供可能日",
    "fixed_date": null
  }}
]

- name：マイルストン名（日本語で簡潔に。上記のラベルを基準に必要に応じて調整してよい）
- fixed_date：YYYY-MM-DD 形式の文字列。日付が読み取れない／書かれていない場合は null を入れる
- 各項目の指定日は最優先する。指定日がない場合でも実行の順番を考慮して項目の順序を決定すること。
  例：報告書の納品日が調査開始日の前になることはないので、調査開始日はかならず報告書納品日の前になる。
  もしすべての項目を実行するのために十分な日程がない場合は、1営業日に複数の項目が入ってもよい。
"""

                    try:
                        response = client.chat.completions.create(
                            model=DEPLOYMENT,
                            messages=[
                                {"role": "system", "content": "あなたは市場調査プロジェクトのPMとして、実務で使えるスケジュール案を作るアシスタントです。"},
                                {"role": "user", "content": prompt},
                            ],
                            temperature=0.4,
                            max_tokens=800,
                        )

                        ai_text = response.choices[0].message.content.strip()

                        # ```json ... ``` で返ってきた場合のガード
                        if ai_text.startswith("```"):
                            ai_text = ai_text.strip("`")
                            ai_text = ai_text.replace("json", "", 1).strip()

                        try:
                            phases = json.loads(ai_text)
                        except Exception:
                            st.error("AI出力をJSONとして解釈できませんでした。出力内容を確認してください。")
                            st.code(ai_text)
                        else:
                            if not isinstance(phases, list):
                                st.error("JSON配列ではありません。出力形式を確認してください。")
                                st.code(ai_text)
                            else:
                                # 中央ペイン（スケジュール案）で利用するためにセッションに保存
                                st.session_state["schedule_phase_draft"] = phases

                                # プレビュー用に DataFrame も保持（任意）
                                try:
                                    import pandas as pd
                                    st.session_state["schedule_phase_draft_df"] = pd.DataFrame(phases)
                                except Exception:
                                    st.session_state["schedule_phase_draft_df"] = None

                                st.success("スケジュールの下書きを作成しました。中央ペインのスケジュール案から参照できるように保存しました。")
                                st.rerun()

                    except Exception as e:
                        st.error(f"AI呼び出しエラー: {e}")

        # ▼ 既に下書きがあればプレビュー表示
        if "schedule_phase_draft" in st.session_state:
            st.markdown("### 抽出されたマイルストン案（プレビュー）")
            try:
                import pandas as pd
                df = st.session_state.get("schedule_phase_draft_df")
                if df is None:
                    df = pd.DataFrame(st.session_state["schedule_phase_draft"])
                st.data_editor(
                    df,
                    hide_index=True,
                    num_rows="fixed",
                    use_container_width=True,
                    key="schedule_phase_draft_preview",
                )
            except Exception:
                # DataFrame化が失敗した場合は生データをそのまま表示
                st.code(st.session_state["schedule_phase_draft"])

    # =========================
    # 右ペイン
    # === 概算見積（仕様入力）===
    elif mode == "概算見積":
        st.subheader("概算見積（仕様入力）")
        st.caption("ここで企画費用（人件費）と実査費用（ベース仕様）を入力すると、中央ペインで5パターンの見積が計算されます。")

        # -------------------------
        # セッション状態の初期値を設定
        # （すでに値があればそのまま維持）
        # -------------------------
        default_values = {
            "hours_plan": 0.0,
            "hours_field": 0.0,
            "hours_agg": 0.0,
            "hours_analysis": 0.0,
            "scr_q": 5,
            "scr_n": 10000,
            "main_q": 20,
            "main_n": 300,
        }
        for k, v in default_values.items():
            if k not in st.session_state:
                st.session_state[k] = v

        # -------------------------
        # ① 企画費用（人件費）
        # -------------------------
        st.markdown("### ① 企画費用（人件費）")

        col1, col2 = st.columns(2)

        with col1:
            st.number_input(
                "調査企画（人時）",
                min_value=0.0,
                step=0.5,
                key="hours_plan",  # ← 中央ペインと同じキー
            )
            st.number_input(
                "調査実査（人時）",
                min_value=0.0,
                step=0.5,
                key="hours_field",
            )

        with col2:
            st.number_input(
                "集計（人時）",
                min_value=0.0,
                step=0.5,
                key="hours_agg",
            )
            st.number_input(
                "分析・報告（人時）",
                min_value=0.0,
                step=0.5,
                key="hours_analysis",
            )

        st.markdown("---")

        # -------------------------
        # ② 実査費用（ベース仕様）
        # -------------------------
        st.markdown("### ② 実査費用（ベース仕様）")

        st.markdown("**スクリーニング調査**")
        cs1, cs2 = st.columns(2)
        with cs1:
            st.number_input(
                "スクリーニング 質問数（問）",
                min_value=0,
                step=1,
                key="scr_q",
            )
        with cs2:
            st.number_input(
                "スクリーニング サンプルサイズ",
                min_value=0,
                step=1000,
                key="scr_n",
            )

        st.markdown("**本調査**")
        cm1, cm2 = st.columns(2)
        with cm1:
            st.number_input(
                "本調査 質問数（問）",
                min_value=0,
                step=1,
                key="main_q",
            )
        with cm2:
            st.number_input(
                "本調査 サンプルサイズ",
                min_value=0,
                step=100,
                key="main_n",
            )

        st.info("※ここで入力した内容をもとに、中央ペインで概算見積（5パターン比較）が自動計算されます。")


    # =========================
    # 右ペイン
    # === PowerPoint出力 ===
    elif mode == "パワーポイントを出力":
        st.subheader("PowerPoint出力")
        st.caption("中央ペインで最終版を作成したあと、ここからダウンロードできます。")

        from pathlib import Path

        pptx_path = st.session_state.get("pptx_path")
        final_path = st.session_state.get("final_pptx_path")

        # final があればそれを、なければ現時点のpptxを候補にする
        candidate_path = pptx_path

        #st.write("DEBUG_pptx_path:", st.session_state.get("pptx_path"))
        #st.write("DEBUG_final_path:", st.session_state.get("final_pptx_path"))


        if candidate_path and Path(candidate_path).is_file():
            label = "📥 最終版PowerPointをダウンロード" if final_path else "📥 現在のPowerPointをダウンロード"

            with open(candidate_path, "rb") as f:
                st.download_button(
                    label,
                    f,
                    file_name=Path(candidate_path).name,
                    use_container_width=True,
                )

            if not final_path:
                st.info("まだ最終版は作成していません。中央ペインの『💾 現在の内容で最終版PowerPointを作成』を押すと、ファイル名付きで確定保存されます。")
        else:
            st.info("中央ペインで最終版を作成すると、ここからダウンロードできるようになります。")

